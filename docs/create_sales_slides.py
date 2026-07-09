"""
売上管理（自費・販売）タブ 操作マニュアル スライド生成スクリプト
実行: python docs/create_sales_slides.py
出力: docs/売上管理タブ_操作マニュアル.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# --- 定数 ---
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT = "Meiryo UI"

C_SELF   = RGBColor(0x7C, 0x3A, 0xED)  # 紫 (自費レンタル)
C_SALE   = RGBColor(0x16, 0xA3, 0x4A)  # 緑 (販売)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_DARK   = RGBColor(0x1F, 0x29, 0x37)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF9, 0xFA, 0xFB)
C_AMBER  = RGBColor(0xD9, 0x77, 0x06)
C_BG     = RGBColor(0xF5, 0xF3, 0xFF)  # 薄紫背景

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # 完全空白


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else C_DARK
    return txb


def slide_header(slide, title, accent=C_SELF):
    # 上部グラデーション風帯
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill=accent)
    add_text(slide, title,
             Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.7),
             size=26, bold=True, color=C_WHITE)
    # 白背景本体
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(6.4), fill=C_LIGHT)


# =====================================================================
# スライド 1 : タイトル
# =====================================================================
sl = prs.slides.add_slide(blank)

# 背景
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xED, 0xE9, 0xFE))

# 上部紫帯
add_rect(sl, 0, 0, SLIDE_W, Inches(0.55), fill=C_SELF)

# 中央コンテンツエリア
add_rect(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.2),
         fill=C_WHITE)

add_text(sl, "売上管理（自費・販売）タブ",
         Inches(1.8), Inches(1.95), Inches(9.8), Inches(1.0),
         size=34, bold=True, color=C_SELF, align=PP_ALIGN.CENTER)
add_text(sl, "操作マニュアル",
         Inches(1.8), Inches(2.95), Inches(9.8), Inches(0.7),
         size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# 色バッジ
badge_y = Inches(3.8)
for bx, bw, bc, bt in [
    (Inches(3.1),  Inches(2.8), C_SELF,  "💰 自費レンタル"),
    (Inches(7.4),  Inches(2.8), C_SALE,  "🛒 販売"),
]:
    add_rect(sl, bx, badge_y, bw, Inches(0.55), fill=bc)
    add_text(sl, bt, bx + Inches(0.1), badge_y + Inches(0.08),
             bw - Inches(0.2), Inches(0.45),
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "WelfareAssist Pro",
         Inches(1.8), Inches(4.85), Inches(9.8), Inches(0.4),
         size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# スライド 2 : 概要
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "売上管理タブとは", accent=C_SELF)

# 左ボックス
add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.8),
         fill=C_WHITE, line=C_SELF)
add_text(sl, "このタブでわかること",
         Inches(0.5), Inches(1.35), Inches(5.6), Inches(0.45),
         size=15, bold=True, color=C_SELF)

items_l = [
    ("💰", "自費レンタル一覧（紫）",  "月額・税込金額・利用期間"),
    ("🛒", "販売一覧（緑）",            "単価・送料・総計・受注/納品日"),
    ("📊", "自動計算",                   "税区分に応じて税込金額を表示"),
    ("🔗", "連動データ",                  "「福祉用具選定」タブと同じデータ"),
]
for i, (ic, t, d) in enumerate(items_l):
    iy = Inches(1.9) + i * Inches(1.2)
    add_text(sl, ic, Inches(0.55), iy, Inches(0.45), Inches(0.45), size=20)
    add_text(sl, t,  Inches(1.05), iy,           Inches(4.9), Inches(0.38),
             size=14, bold=True)
    add_text(sl, d,  Inches(1.05), iy + Inches(0.4), Inches(4.9), Inches(0.38),
             size=11, color=C_GRAY)

# 右ボックス
add_rect(sl, Inches(6.5), Inches(1.3), Inches(6.4), Inches(2.55),
         fill=RGBColor(0xF0, 0xFD, 0xF4), line=C_SALE)
add_text(sl, "✅ このタブは「表示専用」",
         Inches(6.65), Inches(1.38), Inches(6.1), Inches(0.45),
         size=14, bold=True, color=C_SALE)
add_text(sl,
         "データの入力・編集は\n「福祉用具選定」タブで行います。\n\n"
         "このタブは確認・金額チェック用です。",
         Inches(6.65), Inches(1.88), Inches(6.1), Inches(1.6),
         size=12, color=C_DARK)

add_rect(sl, Inches(6.5), Inches(4.05), Inches(6.4), Inches(3.05),
         fill=RGBColor(0xF5, 0xF3, 0xFF), line=C_SELF)
add_text(sl, "📂 タブの表示条件",
         Inches(6.65), Inches(4.13), Inches(6.1), Inches(0.45),
         size=14, bold=True, color=C_SELF)
add_text(sl,
         "• 自費レンタルまたは販売の\n  データがある利用者のみ内容表示\n\n"
         "• 両方ない場合は\n  「データはありません」と表示",
         Inches(6.65), Inches(4.62), Inches(6.1), Inches(2.1),
         size=12, color=C_DARK)

# =====================================================================
# スライド 3 : タブの見方（全体像）
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "タブの画面構成", accent=C_SELF)

# 画面模式図
# 全体フレーム
add_rect(sl, Inches(0.35), Inches(1.25), Inches(12.6), Inches(5.95),
         fill=C_WHITE, line=C_GRAY)

# タブバー
add_rect(sl, Inches(0.35), Inches(1.25), Inches(12.6), Inches(0.5),
         fill=RGBColor(0xE5, 0xE7, 0xEB))
for i, (tl, tc) in enumerate([
    ("基本情報", C_GRAY), ("病歴・状態", C_GRAY),
    ("議事録", C_GRAY), ("変更情報", C_GRAY),
    ("福祉用具選定", C_GRAY), ("売上管理（自費・販売）", C_SELF),
    ("書類管理", C_GRAY),
]):
    tx = Inches(0.5) + i * Inches(1.77)
    bg = RGBColor(0xFF, 0xFF, 0xFF) if tc == C_SELF else RGBColor(0xE5, 0xE7, 0xEB)
    add_rect(sl, tx, Inches(1.3), Inches(1.7), Inches(0.42), fill=bg)
    add_text(sl, tl, tx + Inches(0.03), Inches(1.32), Inches(1.65), Inches(0.38),
             size=7, bold=(tc == C_SELF), color=tc, align=PP_ALIGN.CENTER)

# 自費レンタルセクション
add_rect(sl, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.42),
         fill=C_SELF)
add_text(sl, "💰 自費レンタル  [件数バッジ]",
         Inches(0.6), Inches(1.9), Inches(5.0), Inches(0.35),
         size=11, bold=True, color=C_WHITE)

# テーブルヘッダー（自費）
add_rect(sl, Inches(0.45), Inches(2.27), Inches(12.4), Inches(0.35),
         fill=RGBColor(0xF3, 0xE8, 0xFF))
cols_self = ["商品名", "数量", "月額（税抜）", "税区分", "税込金額", "利用開始日", "利用終了日"]
cw = Inches(12.4) / 7
for ci, ct in enumerate(cols_self):
    add_text(sl, ct, Inches(0.45) + ci * cw, Inches(2.3), cw, Inches(0.3),
             size=8, bold=True, color=C_SELF, align=PP_ALIGN.CENTER)

# テーブルデータ行（自費）
add_rect(sl, Inches(0.45), Inches(2.62), Inches(12.4), Inches(0.3),
         fill=C_WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
sample_self = ["スライディングシート", "1", "¥3,000", "10％", "¥3,300", "2025-04-01", ""]
for ci, ct in enumerate(sample_self):
    add_text(sl, ct, Inches(0.45) + ci * cw, Inches(2.64), cw, Inches(0.26),
             size=8, color=C_DARK if ci not in (4,) else C_SELF,
             bold=(ci == 4), align=PP_ALIGN.CENTER)

# 販売セクション
add_rect(sl, Inches(0.45), Inches(3.05), Inches(12.4), Inches(0.42),
         fill=C_SALE)
add_text(sl, "🛒 販売  [件数バッジ]",
         Inches(0.6), Inches(3.1), Inches(5.0), Inches(0.35),
         size=11, bold=True, color=C_WHITE)

# テーブルヘッダー（販売）
add_rect(sl, Inches(0.45), Inches(3.47), Inches(12.4), Inches(0.35),
         fill=RGBColor(0xDC, 0xFC, 0xE7))
cols_sale = ["商品名", "数量", "単価（税抜）", "税区分", "税込金額", "送料（税抜）", "送料消費税", "総計", "受注日", "納品日", "支払方法", "申請"]
cw2 = Inches(12.4) / 12
for ci, ct in enumerate(cols_sale):
    add_text(sl, ct, Inches(0.45) + ci * cw2, Inches(3.5), cw2, Inches(0.28),
             size=7, bold=True, color=C_SALE, align=PP_ALIGN.CENTER)

# テーブルデータ行（販売）
add_rect(sl, Inches(0.45), Inches(3.82), Inches(12.4), Inches(0.3),
         fill=C_WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
sample_sale = ["歩行器", "1", "¥28,000", "10％", "¥30,800", "¥500", "¥50", "¥31,350", "2025-04-10", "2025-04-15", "現金", "鹿児島市"]
for ci, ct in enumerate(sample_sale):
    add_text(sl, ct, Inches(0.45) + ci * cw2, Inches(3.84), cw2, Inches(0.26),
             size=7, color=C_DARK if ci not in (4, 7) else (C_SELF if ci == 4 else RGBColor(0x16, 0x58, 0x25)),
             bold=(ci in (4, 7)), align=PP_ALIGN.CENTER)

# 凡例
for bx, bc, bt in [
    (Inches(0.5),  C_SELF, "紫 = 自費レンタル"),
    (Inches(3.5),  C_SALE, "緑 = 販売"),
]:
    add_rect(sl, bx, Inches(4.3), Inches(0.28), Inches(0.28), fill=bc)
    add_text(sl, bt, bx + Inches(0.35), Inches(4.3), Inches(2.6), Inches(0.3),
             size=11, bold=True, color=bc)

add_text(sl, "※ 表示のみ。編集・追加は「福祉用具選定」タブで行います",
         Inches(0.5), Inches(4.7), Inches(12.0), Inches(0.35),
         size=12, color=C_AMBER, bold=True)

# =====================================================================
# スライド 4 : 自費レンタルテーブルの見方
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "自費レンタル — テーブルの各列説明", accent=C_SELF)

# 見出し帯
add_rect(sl, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.45),
         fill=C_SELF)
add_text(sl, "💰 自費レンタル（紫テーマ）",
         Inches(0.55), Inches(1.3), Inches(12.0), Inches(0.38),
         size=14, bold=True, color=C_WHITE)

# 列説明テーブル
cols_info = [
    ("商品名",        "福祉用具選定タブで入力した機器名称（請求費目）"),
    ("数量",          "レンタル数量（通常は 1）"),
    ("月額（税抜）",  "1か月あたりの税抜きレンタル料金（unitPrice）"),
    ("税区分",        "「10％」「軽8％」「非課税」のいずれか"),
    ("税込金額",      "月額 × 数量 + 消費税（自動計算）\n  例: ¥3,000 × 1 × 1.10 = ¥3,300"),
    ("利用開始日",    "レンタル開始日（startDate）"),
    ("利用終了日",    "レンタル終了日（endDate）。空欄 = 継続中"),
]

row_h = Inches(0.68)
for i, (col, desc) in enumerate(cols_info):
    ry = Inches(1.8) + i * row_h
    bg = RGBColor(0xF3, 0xE8, 0xFF) if i % 2 == 0 else C_WHITE
    add_rect(sl, Inches(0.4), ry, Inches(12.5), row_h, fill=bg,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_text(sl, col, Inches(0.5), ry + Inches(0.07), Inches(2.8), row_h - Inches(0.1),
             size=13, bold=True, color=C_SELF)
    add_text(sl, desc, Inches(3.4), ry + Inches(0.07), Inches(9.3), row_h - Inches(0.1),
             size=12, color=C_DARK)

# =====================================================================
# スライド 5 : 販売テーブルの見方
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "販売 — テーブルの各列説明", accent=C_SALE)

# 見出し帯
add_rect(sl, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.45),
         fill=C_SALE)
add_text(sl, "🛒 販売（緑テーマ）",
         Inches(0.55), Inches(1.3), Inches(12.0), Inches(0.38),
         size=14, bold=True, color=C_WHITE)

cols_info2 = [
    ("商品名",       "販売した福祉用具の名称"),
    ("数量",         "販売数量"),
    ("単価（税抜）", "1個あたりの税抜き販売価格"),
    ("税区分",       "「10％」「軽8％」「非課税」「税込」"),
    ("税込金額",     "単価 × 数量 + 消費税（自動計算）"),
    ("送料（税抜）", "配送料（税抜）。0円の場合は「-」"),
    ("送料消費税",   "送料 × 10%（自動計算）"),
    ("総計",         "税込金額 + 送料（税抜）+ 送料消費税 + 調整額"),
    ("受注日",       "注文を受けた日付"),
    ("納品日",       "商品を届けた日付。月次売上の基準日"),
    ("支払方法",     "現金・振込・クレジットなど"),
    ("申請",         "福祉用具購入費の申請先市区町村（あり/なし）"),
]

row_h2 = Inches(0.44)
for i, (col, desc) in enumerate(cols_info2):
    ry = Inches(1.8) + i * row_h2
    bg = RGBColor(0xDC, 0xFC, 0xE7) if i % 2 == 0 else C_WHITE
    add_rect(sl, Inches(0.4), ry, Inches(12.5), row_h2, fill=bg,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_text(sl, col, Inches(0.5), ry + Inches(0.04), Inches(2.8), row_h2 - Inches(0.06),
             size=12, bold=True, color=C_SALE)
    add_text(sl, desc, Inches(3.4), ry + Inches(0.04), Inches(9.3), row_h2 - Inches(0.06),
             size=11, color=C_DARK)

# =====================================================================
# スライド 6 : 総計の計算式
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "金額の自動計算ロジック", accent=C_SALE)

# 左列: 自費レンタル計算
add_rect(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(5.9),
         fill=C_WHITE, line=C_SELF)
add_text(sl, "💰 自費レンタル — 税込金額",
         Inches(0.55), Inches(1.32), Inches(5.6), Inches(0.45),
         size=14, bold=True, color=C_SELF)

steps_self = [
    ("①", "小計",      "月額（税抜）× 数量"),
    ("②", "消費税",    "10% の場合: 小計 × 0.10（切捨て）\n8% の場合:  小計 × 0.08（切捨て）\n非課税:      0"),
    ("③", "税込金額",  "小計 + 消費税"),
]
for i, (num, title, formula) in enumerate(steps_self):
    sy = Inches(1.88) + i * Inches(1.6)
    add_rect(sl, Inches(0.5), sy, Inches(0.48), Inches(0.48),
             fill=C_SELF)
    add_text(sl, num, Inches(0.5), sy + Inches(0.05), Inches(0.48), Inches(0.4),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(1.1), sy, Inches(4.8), Inches(0.42),
             size=13, bold=True, color=C_DARK)
    add_text(sl, formula, Inches(1.1), sy + Inches(0.44), Inches(4.8), Inches(1.0),
             size=11, color=C_GRAY)

# 右列: 販売計算
add_rect(sl, Inches(6.7), Inches(1.25), Inches(6.2), Inches(5.9),
         fill=C_WHITE, line=C_SALE)
add_text(sl, "🛒 販売 — 総計の計算式",
         Inches(6.85), Inches(1.32), Inches(5.9), Inches(0.45),
         size=14, bold=True, color=C_SALE)

steps_sale = [
    ("①", "小計",        "単価（税抜）× 数量"),
    ("②", "消費税",      "小計 × 税率（切捨て）\n※「税込」区分は 0"),
    ("③", "税込金額",    "小計 + 消費税"),
    ("④", "送料消費税",  "送料（税抜）× 10%（四捨五入）"),
    ("⑤", "総計",        "税込金額 + 送料（税抜）\n+ 送料消費税 + 調整額"),
]
for i, (num, title, formula) in enumerate(steps_sale):
    sy = Inches(1.88) + i * Inches(1.0)
    add_rect(sl, Inches(6.8), sy, Inches(0.44), Inches(0.44),
             fill=C_SALE)
    add_text(sl, num, Inches(6.8), sy + Inches(0.04), Inches(0.44), Inches(0.38),
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(7.35), sy, Inches(5.4), Inches(0.38),
             size=13, bold=True, color=C_DARK)
    add_text(sl, formula, Inches(7.35), sy + Inches(0.4), Inches(5.4), Inches(0.55),
             size=11, color=C_GRAY)

add_text(sl,
         "⚠ 月次売上処理ページの「販売」サマリーは税抜き総計（送料消費税除く）で集計されます。\n"
         "　 CSV出力の「総計」列は税込み表示のため、数値が異なる場合があります。",
         Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.5),
         size=10, color=C_AMBER, bold=True)

# =====================================================================
# スライド 7 : データ入力の場所（福祉用具選定タブへの案内）
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "データの入力・編集方法", accent=C_SELF)

add_text(sl,
         "売上管理タブは「表示専用」です。データの登録・修正は「福祉用具選定」タブで行います。",
         Inches(0.45), Inches(1.2), Inches(12.4), Inches(0.45),
         size=13, color=C_AMBER, bold=True)

# 自費レンタル入力ガイド
add_rect(sl, Inches(0.4), Inches(1.75), Inches(5.9), Inches(5.45),
         fill=C_WHITE, line=C_SELF)
add_rect(sl, Inches(0.4), Inches(1.75), Inches(5.9), Inches(0.48),
         fill=C_SELF)
add_text(sl, "💰 自費レンタルを登録する",
         Inches(0.55), Inches(1.8), Inches(5.7), Inches(0.38),
         size=14, bold=True, color=C_WHITE)

steps_add_self = [
    ("①", "「福祉用具選定」タブを開く"),
    ("②", "「機器を追加」ボタンをクリック"),
    ("③", "種類:「自費レンタル」を選択"),
    ("④", "属性:「自社物件」または「リース物件」を選択"),
    ("⑤", "商品名・月額・税区分・\n    利用開始日などを入力して保存"),
]
for i, (num, txt) in enumerate(steps_add_self):
    sy = Inches(2.32) + i * Inches(0.95)
    add_rect(sl, Inches(0.55), sy, Inches(0.4), Inches(0.4), fill=C_SELF)
    add_text(sl, num, Inches(0.55), sy + Inches(0.04), Inches(0.4), Inches(0.34),
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, txt, Inches(1.05), sy + Inches(0.04), Inches(5.1), Inches(0.7),
             size=12, color=C_DARK)

# 販売入力ガイド
add_rect(sl, Inches(6.7), Inches(1.75), Inches(6.2), Inches(5.45),
         fill=C_WHITE, line=C_SALE)
add_rect(sl, Inches(6.7), Inches(1.75), Inches(6.2), Inches(0.48),
         fill=C_SALE)
add_text(sl, "🛒 販売を登録する",
         Inches(6.85), Inches(1.8), Inches(5.9), Inches(0.38),
         size=14, bold=True, color=C_WHITE)

steps_add_sale = [
    ("①", "「福祉用具選定」タブを開く"),
    ("②", "「機器を追加」ボタンをクリック"),
    ("③", "種類:「販売」を選択"),
    ("④", "属性:「自社物件」または「リース物件」を選択"),
    ("⑤", "商品名・納品日（必須）・\n    単価・税区分などを入力"),
    ("⑥", "送料・支払方法・申請情報を\n    必要に応じて入力して保存"),
]
for i, (num, txt) in enumerate(steps_add_sale):
    sy = Inches(2.32) + i * Inches(0.82)
    add_rect(sl, Inches(6.85), sy, Inches(0.4), Inches(0.4), fill=C_SALE)
    add_text(sl, num, Inches(6.85), sy + Inches(0.04), Inches(0.4), Inches(0.34),
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, txt, Inches(7.35), sy + Inches(0.04), Inches(5.4), Inches(0.65),
             size=12, color=C_DARK)

# =====================================================================
# スライド 8 : FAQ
# =====================================================================
sl = prs.slides.add_slide(blank)
slide_header(sl, "FAQ — よくある質問", accent=C_SELF)

faqs = [
    ("Q", "「自費レンタル・販売データはありません」と表示される",
     "A", "この利用者に自費レンタルまたは販売の登録がありません。\n"
          "「福祉用具選定」タブで機器を追加してください。"),
    ("Q", "税込金額の計算が合わない",
     "A", "消費税は切り捨て計算です（例: ¥3,001 × 10% → ¥300 切捨 → ¥3,301）。\n"
          "「税込」区分の場合は消費税 0 として計算します。"),
    ("Q", "販売の「総計」と月次売上処理の金額が違う",
     "A", "月次売上処理の販売サマリーは税抜き表示（税込金額 + 送料税抜 + 調整額）、\n"
          "このタブの「総計」は税込み表示（送料消費税も含む）のため差異が生じます。"),
    ("Q", "「申請」列に市区町村名が出ない",
     "A", "福祉用具選定タブの販売フォームで「申請あり」をオン、\n"
          "「申請市区町村」欄に入力されている場合のみ表示されます。"),
    ("Q", "納品日を入れたのに月次売上に出ない",
     "A", "月次売上処理の対象月と納品日の月が一致しているか確認してください。\n"
          "販売は「納品日」の月に計上されます。"),
]

row_h3 = Inches(0.98)
for i, (qt, q, at, a) in enumerate(faqs):
    ry = Inches(1.3) + i * row_h3
    bg = RGBColor(0xF5, 0xF3, 0xFF) if i % 2 == 0 else C_WHITE
    add_rect(sl, Inches(0.4), ry, Inches(12.5), row_h3 - Inches(0.04),
             fill=bg, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(sl, Inches(0.4), ry, Inches(0.4), row_h3 - Inches(0.04), fill=C_SELF)
    add_text(sl, "Q", Inches(0.4), ry + Inches(0.08), Inches(0.4), Inches(0.38),
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, q, Inches(0.88), ry + Inches(0.06), Inches(11.8), Inches(0.38),
             size=12, bold=True, color=C_DARK)
    add_text(sl, a, Inches(0.88), ry + Inches(0.46), Inches(11.8), Inches(0.46),
             size=11, color=C_GRAY)

# =====================================================================
# 保存
# =====================================================================
out = "docs/sales_manual_temp.pptx"
prs.save(out)
# PowerShell rename: Rename-Item docs/sales_manual_temp.pptx "売上管理タブ_操作マニュアル.pptx"
# Final file: docs/売上管理タブ_操作マニュアル.pptx
print("saved: " + out)
