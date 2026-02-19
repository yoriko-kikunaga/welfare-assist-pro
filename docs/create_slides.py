"""
WelfareAssist Pro - 新人向けスライド生成スクリプト
出力: docs/WelfareAssist_Pro_新人向けガイド.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# カラーパレット
COLOR_PRIMARY   = RGBColor(0x2B, 0x6C, 0xB0)  # ダークブルー
COLOR_ACCENT    = RGBColor(0x0D, 0x94, 0x88)  # ティール
COLOR_LIGHT     = RGBColor(0xEF, 0xF6, 0xFF)  # 薄青
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK      = RGBColor(0x1A, 0x20, 0x2C)
COLOR_GRAY      = RGBColor(0x71, 0x7D, 0x8A)
COLOR_SECTION   = {
    'client':       RGBColor(0x2B, 0x6C, 0xB0),
    'kaipoke':      RGBColor(0x6D, 0x28, 0xD9),
    'monthly':      RGBColor(0x6D, 0x28, 0xD9),
    'reconcile':    RGBColor(0x0F, 0x76, 0x6E),
    'receipt':      RGBColor(0xBE, 0x18, 0x5D),
    'ai':           RGBColor(0xB4, 0x5A, 0x09),
    'bed':          RGBColor(0x92, 0x40, 0x00),
    'change':       RGBColor(0x1D, 0x4E, 0xD8),
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank


# ─────────────────────────────────────────────
# ヘルパー関数
# ─────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Meiryo UI"
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  font_size=16, bold=False, color=COLOR_DARK,
                  align=PP_ALIGN.LEFT, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Meiryo UI"
    return txBox


def make_title_bar(slide, bg_color=COLOR_PRIMARY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=bg_color)
    add_rect(slide, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(
        max(0, bg_color[0] - 20), max(0, bg_color[1] - 20), max(0, bg_color[2] - 20)))


def make_content_layout(slide, title_text, title_color=COLOR_PRIMARY):
    """標準コンテンツレイアウト: 上部タイトルバー + 白コンテンツエリア"""
    # 背景白
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_WHITE)
    # タイトルバー
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=title_color)
    # タイトルテキスト
    add_text(slide, title_text, 0.4, 0.15, 12.5, 0.85,
             font_size=28, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    # 下部アクセントライン
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=title_color)
    # ページ番号エリア（右下）
    add_text(slide, "WelfareAssist Pro | 新人向けガイド",
             0.3, 7.1, 12.5, 0.4,
             font_size=10, color=COLOR_WHITE, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# スライド 1: タイトル
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_PRIMARY)
add_rect(slide, 0, 2.8, 13.33, 2.2, fill_color=RGBColor(0x1E, 0x40, 0xAF))
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(0x0D, 0x94, 0x88))

add_text(slide, "WelfareAssist Pro",
         1.5, 1.2, 10, 1.5, font_size=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "新人向けアプリガイド",
         1.5, 2.9, 10, 1.2, font_size=32, bold=False, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)
add_text(slide, "〜福祉用具専門相談員向け業務管理アプリ〜",
         1.5, 4.0, 10, 0.8, font_size=18, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, "株式会社あおぞら",
         1.5, 6.8, 10, 0.6, font_size=14, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# スライド 2: 目次
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "目次", COLOR_PRIMARY)

items = [
    ("01", "アプリの概要", "何のためのアプリか？"),
    ("02", "ログインと画面構成", "最初にすること"),
    ("03", "利用者管理", "6つのタブで情報を管理"),
    ("04", "カイポケCSVインポート", "毎月の介護保険レンタルデータ取込"),
    ("05", "月次売上処理", "売上確認・CSV出力・確定"),
    ("06", "売上・仕入突合", "卸会社請求書との照合"),
    ("07", "レセプトチェック", "請求前の確認チェックリスト"),
    ("08", "AI機能", "議事録生成・OCR・用具提案"),
    ("09", "ベッド管理 / 変更情報", "在庫管理・変更記録の一覧"),
]

col_w = 5.8
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    lx = 0.5 + col * 6.5
    ty = 1.3 + row * 1.12

    # 番号バッジ
    add_rect(slide, lx, ty + 0.05, 0.55, 0.6, fill_color=COLOR_PRIMARY)
    add_text(slide, num, lx, ty + 0.05, 0.55, 0.6,
             font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # タイトル
    add_text(slide, title, lx + 0.65, ty, 5.0, 0.45,
             font_size=16, bold=True, color=COLOR_DARK)
    # 説明
    add_text(slide, desc, lx + 0.65, ty + 0.42, 5.0, 0.45,
             font_size=12, color=COLOR_GRAY)


# ─────────────────────────────────────────────
# スライド 3: アプリ概要
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "01. アプリの概要", COLOR_PRIMARY)

add_text(slide, "WelfareAssist Pro とは",
         0.5, 1.3, 12, 0.6, font_size=22, bold=True, color=COLOR_PRIMARY)
add_text(slide,
         "福祉用具専門相談員の日常業務をひとつのアプリで完結するための業務管理システムです。",
         0.5, 1.85, 12, 0.6, font_size=16, color=COLOR_DARK)

features = [
    ("利用者管理", "8,000件以上の利用者情報をGoogleスプレッドシート＋Kintoneと連携して管理"),
    ("月次売上処理", "介護保険・自費レンタル・販売の売上を月次でCSV出力"),
    ("売上・仕入突合", "卸会社6社の請求書をAI（OCR）で読み取り、粗利を自動計算"),
    ("AI議事録生成", "Google Meetのメモから担当者会議・カンファレンス議事録を自動生成"),
    ("レセプトチェック", "介護保険レンタルの請求前確認チェックリストをアプリ内で管理"),
    ("ベッド管理", "自社所有ベッドの在庫・貸出・消毒・償却を一元管理"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 2.65 + row * 1.3

    add_rect(slide, lx, ty, 6.0, 1.1,
             fill_color=COLOR_LIGHT,
             line_color=RGBColor(0xBF, 0xDB, 0xFE), line_width=0.75)
    add_text(slide, f"✦ {title}", lx + 0.2, ty + 0.08, 5.6, 0.45,
             font_size=15, bold=True, color=COLOR_PRIMARY)
    add_text(slide, desc, lx + 0.2, ty + 0.52, 5.6, 0.5,
             font_size=12, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 4: ログイン・画面構成
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "02. ログインと画面構成", COLOR_PRIMARY)

add_text(slide, "ログイン方法",
         0.5, 1.3, 12, 0.5, font_size=20, bold=True, color=COLOR_PRIMARY)

steps = [
    "ブラウザで https://welfare-assist-pro.web.app を開く",
    "「Googleでサインイン」ボタンをクリック",
    "会社のGoogleアカウント（@社内ドメイン）でログイン",
]
for i, step in enumerate(steps):
    add_rect(slide, 0.5, 1.85 + i * 0.55, 0.45, 0.42, fill_color=COLOR_ACCENT)
    add_text(slide, str(i + 1), 0.5, 1.85 + i * 0.55, 0.45, 0.42,
             font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, step, 1.05, 1.88 + i * 0.55, 11.5, 0.42,
             font_size=15, color=COLOR_DARK)

add_text(slide, "画面構成",
         0.5, 3.75, 12, 0.5, font_size=20, bold=True, color=COLOR_PRIMARY)

layout_items = [
    ("左サイドバー（幅320px）",
     "・ 利用者検索（氏名/カナ/あおぞらID）\n・ 各種ページボタン（福祉用具集計/突合/月次売上/変更情報/ベッド管理/レセプト/ヘルプ）\n・ 利用者一覧（クリックで詳細表示）"),
    ("メインコンテンツ（右側）",
     "・ 利用者詳細（6タブ）または各種専用ページ\n・ 利用者が未選択のときは案内メッセージを表示"),
]
for i, (title, desc) in enumerate(layout_items):
    lx = 0.5 + i * 6.4
    add_rect(slide, lx, 4.3, 6.1, 2.5,
             fill_color=COLOR_LIGHT,
             line_color=RGBColor(0xBF, 0xDB, 0xFE), line_width=0.75)
    add_text(slide, title, lx + 0.2, 4.35, 5.7, 0.5,
             font_size=15, bold=True, color=COLOR_PRIMARY)
    add_multiline(slide, desc.split('\n'), lx + 0.2, 4.85, 5.7, 1.8,
                  font_size=13, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 5: 利用者管理
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "03. 利用者管理", COLOR_SECTION['client'])

add_text(slide, "利用者詳細画面 — 6つのタブ",
         0.5, 1.3, 12, 0.5, font_size=20, bold=True, color=COLOR_SECTION['client'])

tabs = [
    ("Tab 1", "基本情報", "事業所・施設名・ケアマネ・要介護度など基本属性。編集内容はFirestoreに即時保存され、毎日の自動更新後も保持されます。"),
    ("Tab 2", "病歴・医療情報", "病歴・服薬情報を管理。診断書・処方箋などのPDF/画像をAIでテキスト抽出可能。"),
    ("Tab 3", "議事録", "担当者会議・訪問記録を管理。Google Meetのメモ（URL/テキスト/ファイル）からAI議事録を自動生成。"),
    ("Tab 4", "変更情報", "新規・入院・退院・解約などの変更を記録。卸会社への連絡状況も管理。"),
    ("Tab 5", "福祉用具選定", "レンタル・販売品目の一覧・追加。種類→メーカー→商品名の順に選択（マスター連動）。"),
    ("Tab 6", "売上管理", "月別の売上・仕入・粗利を確認。突合済みデータが自動反映。"),
]

for i, (tab, title, desc) in enumerate(tabs):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.95 + row * 1.6

    add_rect(slide, lx, ty, 0.75, 0.55, fill_color=COLOR_SECTION['client'])
    add_text(slide, tab, lx, ty, 0.75, 0.55,
             font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx + 0.85, ty, 5.2, 0.5,
             font_size=15, bold=True, color=COLOR_DARK)
    add_text(slide, desc, lx, ty + 0.6, 6.1, 0.9,
             font_size=12, color=COLOR_GRAY)


# ─────────────────────────────────────────────
# スライド 6: カイポケCSVインポート
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "04. カイポケCSVインポート", COLOR_SECTION['kaipoke'])

add_text(slide, "毎月1回、介護保険レンタルデータをカイポケからインポートします",
         0.5, 1.25, 12, 0.5, font_size=16, color=COLOR_DARK)

add_text(slide, "必要なファイル（2種類）",
         0.5, 1.85, 6, 0.45, font_size=18, bold=True, color=COLOR_SECTION['kaipoke'])
files = [
    ("サービスチェックシート.csv", "カイポケ → 帳票出力 → サービスチェックシート"),
    ("利用者請求.csv",             "カイポケ → 請求管理 → 利用者請求一覧"),
]
for i, (fname, path) in enumerate(files):
    ty = 2.4 + i * 0.7
    add_rect(slide, 0.5, ty, 0.35, 0.55, fill_color=COLOR_SECTION['kaipoke'])
    add_text(slide, str(i+1), 0.5, ty, 0.35, 0.55,
             font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, fname, 0.95, ty, 5.3, 0.35, font_size=14, bold=True, color=COLOR_DARK)
    add_text(slide, path, 0.95, ty + 0.3, 5.3, 0.35, font_size=12, color=COLOR_GRAY)

add_rect(slide, 0.5, 3.85, 5.7, 0.5,
         fill_color=RGBColor(0xFD, 0xF4, 0xFF),
         line_color=RGBColor(0xD8, 0xB4, 0xFE), line_width=1)
add_text(slide, "⚠ 2ファイルとも必ず同じ対象月で出力してください",
         0.6, 3.88, 5.5, 0.45, font_size=13, bold=True, color=RGBColor(0x6D, 0x28, 0xD9))

add_text(slide, "インポート手順",
         6.8, 1.85, 6, 0.45, font_size=18, bold=True, color=COLOR_SECTION['kaipoke'])
steps = [
    "「月次売上処理」→「介護保険レンタル」タブを開く",
    "右上「カイポケCSVインポート」ボタンをクリック",
    "対象月を選択",
    "2つのCSVファイルを順番にアップロード",
    "プレビューで内容を確認",
    "未マッチの利用者は手動で紐づけ",
    "「インポート実行」ボタンをクリック",
]
for i, step in enumerate(steps):
    ty = 2.35 + i * 0.62
    add_rect(slide, 6.8, ty + 0.05, 0.35, 0.42, fill_color=COLOR_SECTION['kaipoke'])
    add_text(slide, str(i + 1), 6.8, ty + 0.05, 0.35, 0.42,
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, step, 7.25, ty + 0.05, 5.6, 0.5, font_size=13, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 7: 月次売上処理
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "05. 月次売上処理", COLOR_SECTION['monthly'])

add_text(slide, "毎月の売上データを確認・確定し、CSVを出力します",
         0.5, 1.25, 12, 0.5, font_size=16, color=COLOR_DARK)

add_text(slide, "3つのタブ",
         0.5, 1.85, 4, 0.45, font_size=18, bold=True, color=COLOR_SECTION['monthly'])

tab_data = [
    ("介護保険\nレンタル", "カイポケCSVインポートで取込んだデータ", RGBColor(0x6D, 0x28, 0xD9)),
    ("自費\nレンタル",     "スプレッドシート連携または手動入力データ", RGBColor(0x7C, 0x3A, 0xED)),
    ("販売",              "スプレッドシート連携または手動入力データ\n利用者負担額・申請額の自動計算", RGBColor(0x8B, 0x5C, 0xF6)),
]
for i, (name, desc, color) in enumerate(tab_data):
    lx = 0.5 + i * 4.1
    add_rect(slide, lx, 2.4, 3.8, 0.8, fill_color=color)
    add_text(slide, name, lx, 2.4, 3.8, 0.8,
             font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx, 3.25, 3.8, 0.9, font_size=12, color=COLOR_DARK)

add_text(slide, "作業フロー",
         0.5, 4.4, 12, 0.45, font_size=18, bold=True, color=COLOR_SECTION['monthly'])

flow = ["月度選択", "事業所フィルター", "データ確認", "CSV出力", "確　定"]
for i, step in enumerate(flow):
    lx = 0.5 + i * 2.5
    add_rect(slide, lx, 5.0, 2.1, 0.7, fill_color=COLOR_SECTION['monthly'])
    add_text(slide, step, lx, 5.0, 2.1, 0.7,
             font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(slide, "▶", lx + 2.15, 5.05, 0.3, 0.6,
                 font_size=18, bold=True, color=COLOR_SECTION['monthly'], align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 5.9, 12.3, 0.75,
         fill_color=RGBColor(0xF5, 0xF3, 0xFF),
         line_color=RGBColor(0xD8, 0xB4, 0xFE), line_width=1)
add_text(slide,
         "💡 確定後はデータが固定されます。修正が必要な場合は「確定解除」ボタンから解除してください。",
         0.7, 5.92, 12.0, 0.7, font_size=13, color=RGBColor(0x6D, 0x28, 0xD9))


# ─────────────────────────────────────────────
# スライド 8: 売上・仕入突合
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "06. 売上・仕入突合", COLOR_SECTION['reconcile'])

add_text(slide, "卸会社からの請求書と社内売上データを照合し、粗利を算出します",
         0.5, 1.25, 12, 0.5, font_size=16, color=COLOR_DARK)

# 突合フロー
flow = ["月度選択", "請求書\nアップロード", "名前\nマッチング", "粗利確認", "CSV出力", "確　定"]
for i, step in enumerate(flow):
    lx = 0.4 + i * 2.12
    add_rect(slide, lx, 1.85, 1.9, 0.8, fill_color=COLOR_SECTION['reconcile'])
    add_text(slide, step, lx, 1.85, 1.9, 0.8,
             font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(slide, "▶", lx + 1.95, 1.93, 0.15, 0.65,
                 font_size=14, color=COLOR_SECTION['reconcile'], align=PP_ALIGN.CENTER)

# 3タブ
add_text(slide, "確認タブ",
         0.5, 2.9, 4, 0.4, font_size=17, bold=True, color=COLOR_SECTION['reconcile'])
tab_data = [
    ("突合済み", "売上と仕入が紐づいた明細\n粗利が自動計算される", RGBColor(0x05, 0x7A, 0x55)),
    ("売上のみ", "仕入データのない売上\n粗利 ＝ 売上額", RGBColor(0x1D, 0x4E, 0xD8)),
    ("仕入のみ", "売上データのない仕入\n→ 要確認・紐づけ対応", RGBColor(0xB4, 0x5A, 0x09)),
]
for i, (name, desc, color) in enumerate(tab_data):
    lx = 0.5 + i * 4.1
    add_rect(slide, lx, 3.4, 3.8, 0.65, fill_color=color)
    add_text(slide, name, lx, 3.4, 3.8, 0.65,
             font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx, 4.1, 3.8, 0.75, font_size=12, color=COLOR_DARK)

# 対応会社
add_text(slide, "対応卸会社（6社）",
         0.5, 5.1, 12, 0.4, font_size=17, bold=True, color=COLOR_SECTION['reconcile'])
companies = ["日建リース工業（PDF）", "野口株式会社（PDF）", "ニシケン（PDF/CSV）",
             "パラマウントケアサービス（PDF/CSV）", "日本ケアサプライ（PDF/CSV）", "株式会社キシヤ（PDF）"]
for i, co in enumerate(companies):
    col = i % 3
    row = i // 3
    lx = 0.5 + col * 4.15
    ty = 5.6 + row * 0.55
    add_rect(slide, lx, ty, 3.9, 0.45,
             fill_color=RGBColor(0xE6, 0xFB, 0xF7),
             line_color=RGBColor(0x6E, 0xE7, 0xB7), line_width=0.75)
    add_text(slide, co, lx + 0.15, ty + 0.04, 3.7, 0.4, font_size=12, color=COLOR_SECTION['reconcile'])


# ─────────────────────────────────────────────
# スライド 9: レセプトチェック
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "07. レセプトチェック", COLOR_SECTION['receipt'])

add_text(slide, "介護保険レンタルの請求前に、請求内容を確認するチェックリスト機能です",
         0.5, 1.25, 12, 0.5, font_size=16, color=COLOR_DARK)

# 手順
add_text(slide, "基本的な流れ",
         0.5, 1.85, 6, 0.45, font_size=18, bold=True, color=COLOR_SECTION['receipt'])
steps = [
    ("月度・事業所を選択", "画面上部で対象月と事業所を選択"),
    ("データ取込", "「データ取込」ボタン → 介護保険レンタル利用者を自動抽出"),
    ("チェックリスト確認", "各利用者のチェック項目をクリックでON/OFF（自動保存）"),
    ("単位数を入力", "直接入力可能。再取込後も手動入力値は保持される"),
    ("月度更新", "解約利用者を除外（amber色ボタン）"),
    ("CSV出力", "請求作業前に必ず出力・保管"),
]
for i, (title, desc) in enumerate(steps):
    ty = 2.4 + i * 0.68
    add_rect(slide, 0.5, ty + 0.07, 0.38, 0.45, fill_color=COLOR_SECTION['receipt'])
    add_text(slide, str(i+1), 0.5, ty + 0.07, 0.38, 0.45,
             font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 0.98, ty, 3.0, 0.38, font_size=14, bold=True, color=COLOR_DARK)
    add_text(slide, desc, 0.98, ty + 0.35, 5.5, 0.38, font_size=12, color=COLOR_GRAY)

# ポイント
add_text(slide, "チェックリストの特徴",
         7.0, 1.85, 6, 0.45, font_size=18, bold=True, color=COLOR_SECTION['receipt'])
points = [
    ("自動取込", "介護保険レンタルがある利用者を自動抽出"),
    ("自動最新化", "ページ読込時に事業所・拠点・入院日等を最新データに自動更新"),
    ("日付自動抽出", "変更情報から入院日・退院日・解約日を自動設定"),
    ("検索対応", "利用者名・居宅介護支援事業所で部分一致検索"),
    ("ソート", "全ヘッダークリックで昇順→降順→解除の3段階"),
    ("事業所フィルター", "切替してもデータは消えない（保存キーは全事業所固定）"),
]
for i, (key, val) in enumerate(points):
    ty = 2.4 + i * 0.68
    add_rect(slide, 7.0, ty + 0.07, 0.9, 0.45, fill_color=COLOR_SECTION['receipt'])
    add_text(slide, key, 7.0, ty + 0.07, 0.9, 0.45,
             font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, val, 8.0, ty + 0.07, 4.8, 0.5, font_size=13, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 10: AI機能
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "08. AI機能（Gemini AI）", COLOR_SECTION['ai'])

add_text(slide, "4つのAI機能が搭載されています",
         0.5, 1.25, 12, 0.45, font_size=16, color=COLOR_DARK)

ai_features = [
    ("AI 議事録生成", "利用者詳細 Tab3",
     "GoogleMeetのメモ（URL/テキスト/ファイル）から\n担当者会議・カンファレンス議事録を自動生成。\n8項目フォーマットで出力。"),
    ("医療文書 OCR", "利用者詳細 Tab2",
     "診断書・処方箋などのPDF/PNG/JPGをアップロードすると\nAIが自動でテキスト抽出（約30秒）。"),
    ("請求書 OCR", "売上・仕入突合",
     "卸会社6社のPDF請求書を自動で読み取り、\n利用者名・商品名・金額を抽出（約30〜60秒）。"),
    ("福祉用具 提案", "利用者詳細 Tab5",
     "利用者の要介護度・病歴・現在の利用品目をもとに\n適切な福祉用具をAIが提案。"),
]

for i, (title, location, desc) in enumerate(ai_features):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.85 + row * 2.45

    add_rect(slide, lx, ty, 6.1, 2.2,
             fill_color=RGBColor(0xFF, 0xF7, 0xED),
             line_color=RGBColor(0xFD, 0xBA, 0x74), line_width=1)
    add_rect(slide, lx, ty, 6.1, 0.55, fill_color=COLOR_SECTION['ai'])
    add_text(slide, title, lx + 0.15, ty, 4.2, 0.55,
             font_size=16, bold=True, color=COLOR_WHITE)
    add_text(slide, location, lx + 4.4, ty, 1.65, 0.55,
             font_size=12, color=RGBColor(0xFE, 0xD7, 0xAA), align=PP_ALIGN.RIGHT)
    add_multiline(slide, desc.split('\n'), lx + 0.15, ty + 0.65, 5.8, 1.5,
                  font_size=13, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 11: ベッド管理 / 変更情報一覧
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "09. ベッド管理 / 変更情報一覧", COLOR_SECTION['bed'])

# 左: ベッド管理
add_rect(slide, 0.3, 1.25, 6.0, 0.55, fill_color=COLOR_SECTION['bed'])
add_text(slide, "自社ベッド管理", 0.3, 1.25, 6.0, 0.55,
         font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

bed_points = [
    ("在庫一覧", "全アイテムの状態（在庫あり/貸出中/消毒中）を管理"),
    ("セット管理", "ベッド＋サイドレール＋マットレスをセットとして一括貸出"),
    ("償却計算",  "購入日・金額・償却月数から月額償却・残存簿価を自動計算"),
    ("消毒記録",  "返却後の消毒開始〜完了を記録（業者名・費用・期間）"),
    ("CSV出力",  "在庫一覧・償却一覧・消毒履歴をCSVで保存"),
]
for i, (key, val) in enumerate(bed_points):
    ty = 2.0 + i * 0.82
    add_rect(slide, 0.3, ty, 1.1, 0.55, fill_color=COLOR_SECTION['bed'])
    add_text(slide, key, 0.3, ty, 1.1, 0.55,
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, val, 1.5, ty + 0.05, 4.7, 0.6, font_size=13, color=COLOR_DARK)

# 右: 変更情報
add_rect(slide, 6.9, 1.25, 6.0, 0.55, fill_color=COLOR_SECTION['change'])
add_text(slide, "変更情報一覧", 6.9, 1.25, 6.0, 0.55,
         font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

change_points = [
    ("確認対象", "新規・入院・退院・解約・変更の記録を全利用者横断で確認"),
    ("フィルター", "期間・事業所・情報種別で絞り込み"),
    ("CSV出力",  "フィルター条件を反映した一覧をCSVで保存"),
    ("スプレッド\nシート", "Googleスプレッドシートへの直接書き出しも可能"),
    ("活用例",   "月初に先月変更分を出力して請求処理・卸会社連絡確認に活用"),
]
for i, (key, val) in enumerate(change_points):
    ty = 2.0 + i * 0.82
    add_rect(slide, 6.9, ty, 1.1, 0.55, fill_color=COLOR_SECTION['change'])
    add_text(slide, key, 6.9, ty, 1.1, 0.55,
             font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, val, 8.1, ty + 0.05, 4.7, 0.65, font_size=13, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 12: 毎月の作業フロー
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "毎月の作業フロー（まとめ）", COLOR_PRIMARY)

add_text(slide, "月初〜月末にかけての標準的な業務フローです",
         0.5, 1.25, 12, 0.45, font_size=16, color=COLOR_DARK)

monthly_flow = [
    ("月初", [
        "カイポケCSVインポート（介護保険レンタル取込）",
        "変更情報一覧で先月分をCSV出力",
    ], RGBColor(0x6D, 0x28, 0xD9)),
    ("月中", [
        "レセプトチェック（請求前確認）",
        "単位数確認・月度更新・CSV出力",
    ], COLOR_SECTION['reconcile']),
    ("月末", [
        "月次売上処理：3タブのデータ確認→CSV出力→確定",
        "売上・仕入突合：請求書アップロード→突合→CSV出力→確定",
    ], COLOR_SECTION['receipt']),
    ("随時", [
        "利用者の変更情報（入院・退院・解約）をTab4に登録",
        "AI議事録生成（担当者会議・訪問後）",
    ], COLOR_SECTION['ai']),
]

for i, (period, items, color) in enumerate(monthly_flow):
    ty = 1.85 + i * 1.3
    add_rect(slide, 0.4, ty, 1.2, 1.05, fill_color=color)
    add_text(slide, period, 0.4, ty, 1.2, 1.05,
             font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, f"・ {item}", 1.8, ty + j * 0.5, 11.0, 0.5,
                 font_size=14, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 13: よくある質問
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
make_content_layout(slide, "よくある質問 / 困ったときは", COLOR_PRIMARY)

faq = [
    ("Q. データが更新されない",
     "毎日 00:00 JSTに自動更新されます。急ぎの場合は手動で再読込してください（ブラウザリロード）。"),
    ("Q. カイポケインポート後にデータが表示されない",
     "月次売上処理で「正しい対象月」が選択されているか確認してください。"),
    ("Q. AI生成が失敗する",
     "処理には30〜60秒かかります。エラーが続く場合はページを再読込してネットワーク接続を確認してください。"),
    ("Q. OCRで利用者が未マッチになる",
     "「未マッチ利用者」の「利用者を選択」から手動で紐づけてください。次回以降は自動マッチングされます。"),
    ("Q. 確定を間違えた",
     "「確定解除」ボタンで解除できます。月次確定がある場合は月次確定解除→各社確定解除の順で行ってください。"),
    ("Q. 使い方がわからない",
     "サイドバー「ヘルプ」ボタンからアプリ内マニュアルを確認できます。"),
]

for i, (q, a) in enumerate(faq):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.3 + row * 1.85

    add_rect(slide, lx, ty, 6.1, 1.65,
             fill_color=COLOR_LIGHT,
             line_color=RGBColor(0xBF, 0xDB, 0xFE), line_width=0.75)
    add_rect(slide, lx, ty, 6.1, 0.5, fill_color=COLOR_PRIMARY)
    add_text(slide, q, lx + 0.15, ty, 5.9, 0.5,
             font_size=13, bold=True, color=COLOR_WHITE)
    add_multiline(slide, [a], lx + 0.15, ty + 0.58, 5.8, 0.95,
                  font_size=12, color=COLOR_DARK)


# ─────────────────────────────────────────────
# スライド 14: まとめ
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_PRIMARY)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill_color=RGBColor(0x0D, 0x94, 0x88))

add_text(slide, "WelfareAssist Pro",
         1.5, 1.0, 10, 1.2, font_size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "ぜひ積極的に活用してください！",
         1.5, 2.2, 10, 1.0, font_size=28, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

summary_items = [
    "利用者情報の一元管理",
    "毎月の業務をアプリ内で完結",
    "AIで議事録・OCRを自動化",
    "ヘルプボタンでいつでも確認可",
]
for i, item in enumerate(summary_items):
    lx = 1.5 + (i % 2) * 5.5
    ty = 3.3 + (i // 2) * 0.75
    add_text(slide, f"✓  {item}", lx, ty, 5.2, 0.65,
             font_size=17, color=COLOR_WHITE)

add_text(slide, "https://welfare-assist-pro.web.app",
         1.5, 5.6, 10, 0.7, font_size=18, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "わからないことは先輩スタッフや管理者に気軽に相談してください",
         1.5, 6.2, 10, 0.6, font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# 保存
# ─────────────────────────────────────────────
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "WelfareAssist_Pro_新人向けガイド.pptx")
prs.save(out_path)
print("Saved: " + out_path)
