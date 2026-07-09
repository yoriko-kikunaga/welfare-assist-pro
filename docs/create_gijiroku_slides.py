"""
WelfareAssist Pro - 議事録タブ 操作マニュアル（スライド生成）
出力: docs/議事録タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_gijiroku_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（他タブマニュアルと共通）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)   # ティール: メインアクセント
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 議事録タブ固有色
C_AI       = RGBColor(0x7C, 0x3A, 0xED)   # 紫: AI機能
C_AILT     = RGBColor(0xED, 0xE9, 0xFE)
C_MEET     = RGBColor(0x16, 0xA3, 0x4A)   # 緑: Meetメモ取込
C_MEETLT   = RGBColor(0xDC, 0xFC, 0xE7)
C_SHEET    = RGBColor(0x1D, 0x4E, 0xD8)   # 青: スプレッドシート同期
C_SHEETLT  = RGBColor(0xDB, 0xEA, 0xFE)
C_CONF     = RGBColor(0xD9, 0x77, 0x06)   # amber: カンファレンス
C_CONFLT   = RGBColor(0xFF, 0xF7, 0xED)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT, footer="議事録タブ 操作マニュアル"):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro  |  " + footer,
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def callout(slide, l, t, w, h, title, body, warn=False, ai=False, meet=False, sheet=False):
    if ai:
        bg, edge = C_AILT, C_AI
    elif meet:
        bg, edge = C_MEETLT, C_MEET
    elif sheet:
        bg, edge = C_SHEETLT, C_SHEET
    elif warn:
        bg, edge = C_WARNLT, C_WARN
    else:
        bg, edge = C_ACCENTLT, C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def meeting_card_header(slide, y, badge, badge_bg, badge_fg, date, recorder, h=0.6):
    """議事録カードのヘッダー行"""
    add_rect(slide, 0.7, y, 11.9, h, fill=badge_bg)
    # 種別バッジ
    add_rect(slide, 0.82, y + 0.1, 1.6, h - 0.2, fill=C_WHITE)
    add_text(slide, badge, 0.84, y + 0.1, 1.56, h - 0.2,
             size=11, bold=True, color=badge_fg, anchor=MSO_ANCHOR.MIDDLE)
    # 日付
    add_text(slide, date, 2.58, y + 0.06, 2.5, h - 0.12, size=13, bold=True, color=C_INK)
    # 記録者
    add_text(slide, "記録者: " + recorder, 5.2, y + 0.06, 3.5, h - 0.12, size=12, color=C_MUTED)


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5,
         size=16, bold=True, color=C_ACCENT)
add_text(s, "「議事録」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=46, bold=True, color=C_INK)
add_rect(s, 1.05, 3.95, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("担当者会議・カンファレンスの記録を入力・管理する画面です。", False, C_MUTED),
    ("Meetメモの取込・AIによる議事録生成・スプレッドシートへの自動転記にも対応しています。", False, C_MUTED),
], 1.05, 4.25, 11, 1.0, size=17)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：利用者をクリック → 「議事録」タブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. このタブでできること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")

feats = [
    (C_ACCENTLT, C_ACCENT, "📝", "議事録を記録する",
     "担当者会議・カンファレンスの日時・出席者・内容を記録。出席者・CM・SWなど詳細なフィールドを入力できます。"),
    (C_MEETLT, C_MEET, "💬", "Meetメモから作成する",
     "Google Meet の議事メモ（URL・テキスト・PDF/TXT）を取り込み、そのまま議事録の素材にできます。"),
    (C_AILT, C_AI, "🤖", "AIで議事録を生成する",
     "内容を入力後「AI議事録を生成」を押すと、会議種別に合わせた構造化サマリーをAIが自動作成します。"),
    (C_SHEETLT, C_SHEET, "📊", "スプレッドシートへ自動転記",
     "保存後4秒以内に Googleスプレッドシートの「議事録」シートへ自動で追記されます（重複なし）。"),
]
x = 0.7
for bg, edge, ic, title, body in feats:
    add_rect(s, x, 2.0, 2.8, 3.4, fill=bg, line=edge, lw=1)
    add_rect(s, x, 2.0, 0.1, 3.4, fill=edge)
    add_text(s, ic, x + 0.25, 2.15, 0.8, 0.7, size=26, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x + 0.22, 2.9, 2.4, 0.55, size=15, bold=True, color=C_INK)
    add_text(s, body, x + 0.22, 3.5, 2.45, 1.65, size=12, color=C_MUTED)
    x += 2.98

callout(s, 0.7, 5.7, 11.9, 0.95,
        "自動保存について",
        "入力内容は1.2秒後に自動的に保存されます。「保存」ボタンを押さなくても大丈夫です。")

# ════════════════════════════════════════════════
# 3. 議事録一覧の見方
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "議事録一覧の見方", eyebrow="LIST VIEW")

# カード1: カンファレンス（展開状態）
y = 1.85
add_rect(s, 0.7, y, 11.9, 0.58, fill=C_CONFLT)
add_rect(s, 0.82, y + 0.1, 1.7, 0.38, fill=C_WHITE)
add_text(s, "カンファレンス時", 0.84, y + 0.1, 1.66, 0.38,
         size=10, bold=True, color=C_CONF, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "2026-07-01", 2.68, y + 0.07, 2.2, 0.42, size=13, bold=True, color=C_INK)
add_text(s, "記録者: 菊永　　場所: 特養あおぞら", 5.0, y + 0.07, 5.5, 0.42, size=12, color=C_MUTED)
# 展開アイコン（上向き矢印）
add_rect(s, 11.8, y + 0.12, 0.5, 0.3, fill=C_ACCENT)
add_text(s, "▲", 11.8, y + 0.12, 0.5, 0.3, size=10, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# カード2: 担当者会議（折りたたみ状態）
y2 = y + 0.72
add_rect(s, 0.7, y2, 11.9, 0.58, fill=C_ACCENTLT)
add_rect(s, 0.82, y2 + 0.1, 1.7, 0.38, fill=C_WHITE)
add_text(s, "担当者会議（新規）", 0.84, y2 + 0.1, 1.66, 0.38,
         size=10, bold=True, color=C_ACCENT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "2026-06-15", 2.68, y2 + 0.07, 2.2, 0.42, size=13, bold=True, color=C_INK)
add_text(s, "記録者: 田中　　居宅: ○○介護支援センター", 5.0, y2 + 0.07, 5.5, 0.42, size=12, color=C_MUTED)
add_rect(s, 11.8, y2 + 0.12, 0.5, 0.3, fill=C_LINE)
add_text(s, "▼", 11.8, y2 + 0.12, 0.5, 0.3, size=10, bold=True, color=C_MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_table(s, [
    ["要素", "説明"],
    ["ヘッダー色（ティール）", "担当者会議（新規・更新・退院時）・カンファレンス時"],
    ["ヘッダー色（橙）", "種別が「その他」の場合のみ"],
    ["▲ で展開 / ▼ で折りたたみ", "カードをクリックして詳細を表示・非表示"],
    ["追加した順（上が最新）", "記録を追加した順に上から表示されます"],
], 0.7, 4.45, 11.9, 1.8, [3.0, 8.9], size=13)

# ════════════════════════════════════════════════
# 4. 議事録を新規追加する
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "議事録を新規追加する", eyebrow="ADD MEETING")

steps = [
    ("1", "「記録を追加」を押す",
     "画面右上の緑のボタンを押します。新しいカードが最上部に追加され、編集モードになります。"),
    ("2", "基本情報を入力する",
     "日付・種別・記録者・場所などを入力します。事業所は基本情報タブの設定が自動で入ります。"),
    ("3", "内容・詳細を記入して保存",
     "議事録内容（テキストエリア）を入力します。1.2秒後に自動保存されます。"),
]
x = 0.7
for n, h, b in steps:
    add_rect(s, x, 2.1, 3.75, 2.6, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, n, x + 0.25, 2.25, 1.2, 1.0, size=44, bold=True, color=C_ACCENT)
    add_text(s, h, x + 0.25, 3.25, 3.3, 0.55, size=18, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.85, 3.3, 0.95, size=13, color=C_MUTED)
    x += 4.05

callout(s, 0.7, 5.05, 11.9, 0.88,
        "「Meetメモから作成」ボタン（緑）との違い",
        "「記録を追加」は空のフォームから手動入力するときに使います。Google Meet 等の文字起こしがある場合は「Meetメモから作成」が便利です。",
        meet=True)

callout(s, 0.7, 6.1, 11.9, 0.7,
        "居宅介護支援事業所・担当CMは基本情報タブから自動入力されます",
        "これらのフィールドは読み取り専用です。変更したい場合は「基本情報」タブで編集してください。")

# ════════════════════════════════════════════════
# 5. 入力フィールド一覧
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "入力フィールド一覧", eyebrow="FIELDS")

add_table(s, [
    ["フィールド", "内容", "備考"],
    ["日付", "会議実施日（YYYY-MM-DD）", ""],
    ["種別", "カンファレンス時 / 担当者会議（新規・更新・退院時）/ その他", "選択式"],
    ["事業所", "鹿児島（ACG）/ 福岡（Lichi）", "基本情報タブから自動入力・変更不可"],
    ["記録者", "担当者名（自由入力）", ""],
    ["施設名", "開催場所の施設名", ""],
    ["出席者", "会議の参加者（自由入力）", ""],
    ["居宅介護支援事業所", "担当ケアマネの所属事業所", ""],
    ["担当CM", "担当ケアマネジャー名", ""],
    ["病院名", "関連する病院名（入院・退院時等）", ""],
    ["担当SW", "担当ソーシャルワーカー名", ""],
    ["利用区分", "介護保険レンタル / 自費レンタル / 購入 / 併用", "選択式"],
    ["ケアプラン", "確認済 / 未確認", "選択式"],
    ["提供票", "確認済 / 未確認", "選択式"],
    ["リマインダー", "あり / なし", "選択式"],
    ["議事録内容", "会議内容の本文（自由テキスト）", "AI生成の入力元"],
    ["AIサマリー", "AIが生成した構造化議事録", "「AI議事録を生成」ボタンで作成"],
], 0.7, 1.85, 11.9, 5.25, [2.4, 6.3, 3.2], size=11)

# ════════════════════════════════════════════════
# 6. Meetメモから作成する
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "Meetメモ・外部テキストから作成する", eyebrow="MEET IMPORT", bar=C_MEET)

add_rect(s, 0.7, 2.0, 11.9, 1.2, fill=C_MEETLT)
add_rect(s, 0.7, 2.0, 0.1, 1.2, fill=C_MEET)
add_lines(s, [
    ("「Meetメモから作成」ボタンを押すと、Google Meet の文字起こしや会議メモを取り込むモーダルが開きます。", False, C_INK),
    ("取り込まれたテキストは「議事録内容」に自動セットされ、AI議事録生成も自動で始まります。", False, C_INK),
], 0.95, 2.1, 11.2, 1.0, size=15)

methods = [
    (C_MEETLT, C_MEET, "URL取込",
     "Google ドキュメントのURLを貼り付けます。\n「リンクを知っている全員に公開」設定のドキュメントに対応。"),
    (C_ACCENTLT, C_ACCENT, "テキスト貼り付け",
     "Google Meet の「文字起こし」テキストを\nコピーしてそのまま貼り付けます。"),
    (C_LINE, C_MUTED, "ファイルアップロード",
     "会議メモの .txt または .pdf ファイルを\nアップロードします（PDF はAIがテキスト抽出）。"),
]
x = 0.7
for bg, edge, title, body in methods:
    add_rect(s, x, 3.45, 3.8, 2.3, fill=bg, line=edge, lw=1.2)
    add_rect(s, x, 3.45, 0.09, 2.3, fill=edge)
    add_text(s, title, x + 0.25, 3.6, 3.3, 0.55, size=18, bold=True, color=C_INK)
    add_text(s, body, x + 0.25, 4.2, 3.35, 1.35, size=13, color=C_MUTED)
    x += 4.07

callout(s, 0.7, 6.0, 11.9, 0.75,
        "取込後はAIサマリー生成が自動でスタートします",
        "「Meetメモから作成」で議事録を追加すると、AIサマリー生成が自動的に開始されます。完了まで数秒お待ちください。",
        meet=True)

# ════════════════════════════════════════════════
# 7. AI議事録を生成する
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "AI議事録を生成する", eyebrow="AI SUMMARY", bar=C_AI)

add_rect(s, 0.7, 2.0, 11.9, 1.1, fill=C_AILT)
add_rect(s, 0.7, 2.0, 0.1, 1.1, fill=C_AI)
add_lines(s, [
    ("「議事録内容」テキストエリアに会議の内容を入力後、「AI議事録を生成」ボタン（紫）を押すとAIが内容を構造化します。", False, C_INK),
    ("会議種別によって生成フォーマットが変わります（担当者会議・カンファレンス ／ その他）。", False, C_INK),
], 0.95, 2.1, 11.2, 0.9, size=14)

formats = [
    (C_AILT, C_AI,
     "担当者会議 / カンファレンス時",
     ["会議目的", "出席者・所属", "利用者の現状", "協議内容", "決定事項", "今後の対応・役割分担", "次回予定", "特記事項"],
     8),
    (C_ACCENTLT, C_ACCENT,
     "その他（訪問・電話対応等）",
     ["訪問日時", "訪問目的", "利用者の状態", "確認事項", "対応内容", "今後の予定", "特記事項"],
     7),
]
x = 0.7
for bg, edge, title, items, cnt in formats:
    add_rect(s, x, 3.28, 5.7, 3.05, fill=bg, line=edge, lw=1)
    add_rect(s, x, 3.28, 0.09, 3.05, fill=edge)
    add_text(s, title, x + 0.25, 3.4, 5.2, 0.5, size=15, bold=True, color=C_INK)
    add_text(s, str(cnt) + "項目", x + 0.25, 3.85, 5.2, 0.4, size=22, bold=True, color=edge)
    bullets = "  ".join(["・" + itm for itm in items])
    add_text(s, bullets, x + 0.25, 4.3, 5.2, 1.75, size=12, color=C_MUTED)
    x += 6.0

callout(s, 0.7, 6.45, 11.9, 0.72,
        "生成結果は「AIサマリー」欄に表示されます",
        "「議事録内容」の下にある「AIサマリー」テキストエリアに結果が入ります。内容を確認・修正してください。",
        ai=True)

# ════════════════════════════════════════════════
# 8. スプレッドシートへの自動転記
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "スプレッドシートへの自動転記", eyebrow="SPREADSHEET SYNC", bar=C_SHEET)

add_rect(s, 0.7, 2.0, 11.9, 1.05, fill=C_SHEETLT)
add_rect(s, 0.7, 2.0, 0.1, 1.05, fill=C_SHEET)
add_lines(s, [
    ("議事録を保存すると、約4秒後に自動的に Google スプレッドシートの「議事録」シートへ追記されます。", False, C_INK),
    ("同じ記録が二重に追記されることはありません（レコードIDで重複管理）。", False, C_INK),
], 0.95, 2.1, 11.2, 0.88, size=14)

add_table(s, [
    ["シート列", "内容"],
    ["レコードID", "システムが自動発行する固有ID（変更不可）"],
    ["日付", "会議実施日"],
    ["あおぞらID / 利用者名 / 事業所", "利用者情報（自動）"],
    ["種別 / 記録者 / 施設名 / 出席者", "議事録の基本情報"],
    ["居宅事業所 / 担当CM / 病院名 / 担当SW", "関係者情報"],
    ["利用区分 / ケアプラン / 提供票 / リマインダー", "確認事項"],
    ["議事録内容 / AIサマリー", "内容フィールド（最大5,000 / 2,000文字）"],
], 0.7, 3.25, 11.9, 3.5, [3.0, 8.9], size=13)

callout(s, 0.7, 6.88, 11.9, 0.72,
        "毎日深夜にも自動一括転記されます",
        "アプリ保存時の即時転記に加え、毎日0時（JST）の自動バッチでも全員分の未転記レコードが一括で追記されます。",
        sheet=True)

# ════════════════════════════════════════════════
# 9. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="FAQ")

qas = [
    ("内容を誤って入力・上書きしてしまいました",
     "自動保存されてしまった場合、スプレッドシートの「議事録」シートに保存前の内容が残っていることがあります。確認してみてください。"),
    ("AIサマリーが生成されません",
     "「議事録内容」テキストエリアが空の場合は生成されません。先に内容を入力してから「AI議事録を生成」を押してください。"),
    ("Meetメモ取込でURLを入れたがエラーになります",
     "Google ドキュメントが「リンクを知っている全員に公開（閲覧可）」になっているか確認してください。"),
    ("スプレッドシートに転記されていません",
     "保存後4〜5秒お待ちください。それでも反映されない場合、翌日の深夜バッチで自動転記されます。"),
    ("同じ議事録が2行追記されていますか？",
     "レコードIDで重複チェックしているため通常は起きません。重複している場合はスプレッドシート上で1行削除してください。"),
]
y = 2.0
rh = 0.96
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_ACCENT)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_ACCENT)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.38, size=14, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.44, 10.9, 0.44, size=12, color=C_MUTED)
    y += rh + 0.08

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "議事録タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
