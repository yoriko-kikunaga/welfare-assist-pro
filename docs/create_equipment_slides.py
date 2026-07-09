"""
WelfareAssist Pro - 福祉用具選定タブ 操作マニュアル（スライド生成）
出力: docs/福祉用具選定タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_equipment_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（他タブマニュアルと共通）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 福祉用具選定タブ 3種別の色
C_INS      = RGBColor(0x1D, 0x4E, 0xD8)   # 青: 介護保険レンタル
C_INSLT    = RGBColor(0xDB, 0xEA, 0xFE)
C_SELF     = RGBColor(0x7C, 0x3A, 0xED)   # 紫: 自費レンタル
C_SELFLT   = RGBColor(0xED, 0xE9, 0xFE)
C_SALE     = RGBColor(0x16, 0xA3, 0x4A)   # 緑: 販売
C_SALELT   = RGBColor(0xDC, 0xFC, 0xE7)
C_ORANGE   = RGBColor(0xC2, 0x57, 0x0E)   # 橙: 自社物件
C_ORANGELT = RGBColor(0xFF, 0xED, 0xD5)
C_TEAL     = C_ACCENT
C_TEALLT   = C_ACCENTLT

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro  |  福祉用具選定タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def callout(slide, l, t, w, h, title, body, warn=False, ins=False, sale=False, self_=False):
    if ins:
        bg, edge = C_INSLT, C_INS
    elif self_:
        bg, edge = C_SELFLT, C_SELF
    elif sale:
        bg, edge = C_SALELT, C_SALE
    elif warn:
        bg, edge = C_WARNLT, C_WARN
    else:
        bg, edge = C_ACCENTLT, C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=12):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def section_header(slide, y, label, count_label, color, colorlt):
    """セクションヘッダーバー"""
    add_rect(slide, 0.7, y, 11.9, 0.5, fill=color)
    add_text(slide, label, 0.95, y + 0.06, 7.0, 0.38,
             size=16, bold=True, color=C_WHITE)
    add_rect(slide, 8.5, y + 0.1, 1.4, 0.3, fill=colorlt)
    add_text(slide, count_label, 8.52, y + 0.1, 1.36, 0.3,
             size=11, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def table_row(slide, y, cells, cell_w, bg, rh=0.38):
    """テーブル行"""
    x = 0.7
    for text, w in zip(cells, cell_w):
        add_rect(slide, x, y, w, rh, fill=bg, line=C_LINE, lw=0.3)
        add_text(slide, text, x + 0.06, y + 0.04, w - 0.12, rh - 0.08,
                 size=11, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
        x += w


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5,
         size=16, bold=True, color=C_ACCENT)
add_text(s, "「福祉用具選定」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=44, bold=True, color=C_INK)
add_rect(s, 1.05, 3.95, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("利用者に提供している福祉用具を登録・管理する画面です。", False, C_MUTED),
    ("介護保険レンタル・自費レンタル・販売の3種類を登録できます。", False, C_MUTED),
    ("機器マスターから商品を選択すると、タイスコードなどが自動入力されます。", False, C_MUTED),
], 1.05, 4.25, 11, 1.2, size=17)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：利用者をクリック → 「福祉用具選定」タブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. このタブでできること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")

feats = [
    (C_INSLT, C_INS, "🏥", "介護保険レンタルを登録",
     "介護保険適用のレンタル機器を登録します。カイポケCSV取込後は自動で反映されます。"),
    (C_SELFLT, C_SELF, "💰", "自費レンタルを登録",
     "自費契約のレンタル機器を登録します。単価・数量・利用期間を入力して管理します。"),
    (C_SALELT, C_SALE, "🛒", "販売を登録",
     "福祉用具の販売を記録します。受注日・納品日・支払方法・申請状況を管理できます。"),
    (C_ACCENTLT, C_ACCENT, "🔍", "マスターから商品選択",
     "機器マスターから種類→メーカー→商品名の順に絞り込み選択できます。タイスコードが自動入力されます。"),
]
x = 0.7
for bg, edge, ic, title, body in feats:
    add_rect(s, x, 2.0, 2.8, 3.4, fill=bg, line=edge, lw=1)
    add_rect(s, x, 2.0, 0.1, 3.4, fill=edge)
    add_text(s, ic, x + 0.25, 2.15, 0.8, 0.7, size=26, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x + 0.22, 2.9, 2.45, 0.55, size=15, bold=True, color=C_INK)
    add_text(s, body, x + 0.22, 3.5, 2.45, 1.65, size=12, color=C_MUTED)
    x += 2.98

callout(s, 0.7, 5.7, 11.9, 0.95,
        "介護保険レンタルはカイポケCSVで自動更新されます",
        "カイポケCSVをインポートすると介護保険レンタルの機器一覧が自動で最新化されます。自費レンタル・販売は手動登録です。")

# ════════════════════════════════════════════════
# 3. 一覧画面の見方（3種別）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "一覧画面の見方", eyebrow="LIST VIEW")

# 介護保険レンタルセクション
section_header(s, 1.85, "🏥  介護保険レンタル", "3件", C_INS, C_INSLT)
headers_ins = ["商品名", "メーカー", "卸会社", "種類", "単位数", "利用開始日", "利用終了日", "カイポケ"]
cw_ins = [2.5, 1.5, 1.4, 1.4, 0.8, 1.0, 1.0, 1.0]
y = 2.42
add_rect(s, 0.7, y, 11.9, 0.38, fill=C_INSLT)
x = 0.7
for h, w in zip(headers_ins, cw_ins):
    add_text(s, h, x + 0.05, y + 0.06, w - 0.1, 0.26, size=10, bold=True, color=C_INS)
    x += w

rows_ins = [
    ("スーパーショートベッド", "パラマウント", "パラマウントCS", "特殊寝台", "52", "2026-04-01", "", "登録済"),
    ("介助バー（縦手すり付）", "パラマウント", "パラマウントCS", "特殊寝台付属品", "36", "2026-04-01", "", "登録済"),
]
for i, row in enumerate(rows_ins):
    bg = C_WHITE if i % 2 == 0 else C_INSLT
    table_row(s, y + 0.38 + i * 0.38, row, cw_ins, bg)

# 自費レンタルセクション
section_header(s, 3.6, "💰  自費レンタル", "1件", C_SELF, C_SELFLT)
y2 = 4.17
add_rect(s, 0.7, y2, 11.9, 0.38, fill=C_SELFLT)
for j, (h, w) in enumerate([("商品名", 3.2), ("卸会社", 1.8), ("単価", 1.2), ("数量", 0.8), ("税込金額", 1.4), ("利用開始日", 1.5), ("利用終了日", 1.5)]):
    x2 = 0.7 + sum([3.2,1.8,1.2,0.8,1.4,1.5,1.5][:j])
    add_text(s, h, x2 + 0.05, y2 + 0.06, w - 0.1, 0.26, size=10, bold=True, color=C_SELF)
table_row(s, y2 + 0.38, ("防水シーツ（大）", "野口", "¥1,200", "1", "¥1,320", "2026-05-01", ""), [3.2,1.8,1.2,0.8,1.4,1.5,1.5], C_WHITE)

add_table(s, [
    ["表示・操作", "意味"],
    ["ヘッダーの色（青/紫/緑）", "介護保険レンタル / 自費レンタル / 販売 を区別"],
    ["行クリック（編集モード時）", "その機器の詳細入力フォームが開きます"],
    ["🔒マーク", "月次売上が確定済みのため編集・削除できません"],
    ["削除ボタン（赤）", "その機器を一覧から削除します（元に戻せません）"],
], 0.7, 5.62, 11.9, 1.65, [3.0, 8.9], header_color=C_ACCENT)

# ════════════════════════════════════════════════
# 4. 機器を追加する（2ステップモーダル）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "機器を追加する", eyebrow="ADD EQUIPMENT")

# ステップ1
add_rect(s, 0.7, 2.0, 5.7, 4.15, fill=C_WHITE, line=C_LINE, lw=1)
add_rect(s, 0.7, 2.0, 5.7, 0.45, fill=C_ACCENTLT)
add_rect(s, 0.7, 2.0, 0.09, 0.45, fill=C_ACCENT)
add_text(s, "Step 1  種類を選択", 0.95, 2.07, 5.2, 0.32, size=14, bold=True, color=C_ACCENT)

step1_items = [
    (C_INSLT, C_INS, "🏥", "介護保険レンタル", "介護保険適用のレンタル用具"),
    (C_SELFLT, C_SELF, "💰", "自費レンタル", "自費でのレンタル用具"),
    (C_SALELT, C_SALE, "🛒", "販売", "福祉用具の販売"),
]
y = 2.56
for bg, edge, ic, label, sub in step1_items:
    add_rect(s, 0.85, y, 5.35, 0.75, fill=bg, line=edge, lw=0.8)
    add_text(s, ic, 0.98, y + 0.1, 0.5, 0.55, size=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, label, 1.55, y + 0.1, 4.4, 0.3, size=14, bold=True, color=edge)
    add_text(s, sub, 1.55, y + 0.4, 4.4, 0.25, size=11, color=C_MUTED)
    y += 0.85

# ステップ2
add_rect(s, 7.0, 2.0, 5.6, 4.15, fill=C_WHITE, line=C_LINE, lw=1)
add_rect(s, 7.0, 2.0, 5.6, 0.45, fill=C_ORANGELT)
add_rect(s, 7.0, 2.0, 0.09, 0.45, fill=C_ORANGE)
add_text(s, "Step 2  属性を選択", 7.25, 2.07, 5.1, 0.32, size=14, bold=True, color=C_ORANGE)
add_text(s, "種類: 介護保険レンタル", 7.25, 2.52, 5.1, 0.3, size=12, bold=True, color=C_INS)

step2_items = [
    (C_ORANGELT, C_ORANGE, "🏠", "自社物件", "自社所有の福祉用具（仕入なし）"),
    (C_TEALLT, C_TEAL, "📋", "リース物件", "リース契約の福祉用具"),
]
y2 = 2.9
for bg, edge, ic, label, sub in step2_items:
    add_rect(s, 7.15, y2, 5.25, 0.75, fill=bg, line=edge, lw=0.8)
    add_text(s, ic, 7.28, y2 + 0.1, 0.5, 0.55, size=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, label, 7.85, y2 + 0.1, 4.3, 0.3, size=14, bold=True, color=edge)
    add_text(s, sub, 7.85, y2 + 0.4, 4.3, 0.25, size=11, color=C_MUTED)
    y2 += 0.85

add_text(s, "← 種類選択に戻る", 7.25, 4.6, 5.1, 0.3, size=11, color=C_MUTED)

# 矢印
add_rect(s, 6.25, 3.8, 0.7, 0.35, fill=C_ACCENT)
add_text(s, "→", 6.25, 3.8, 0.7, 0.35, size=16, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

callout(s, 0.7, 6.35, 11.9, 0.82,
        "Step 2 の後、機器詳細の入力フォームが自動で開きます",
        "種類と属性を選択するとすぐに詳細入力モーダルが表示されます。内容を入力して「保存」を押してください。")

# ════════════════════════════════════════════════
# 5. 介護保険レンタルの入力フォーム
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "介護保険レンタルの入力フォーム", eyebrow="INSURANCE RENTAL FORM", bar=C_INS)

# カスケードフィルタ説明
add_rect(s, 0.7, 2.0, 11.9, 1.05, fill=C_INSLT)
add_rect(s, 0.7, 2.0, 0.09, 1.05, fill=C_INS)
add_lines(s, [
    ("種類 → メーカー → 商品名 の順に選択すると、タイスコード・単位数が自動入力されます。", True, C_INS),
    ("この「カスケード選択」を使うと入力ミスを防げます。手入力も可能です。", False, C_INK),
], 0.95, 2.1, 11.2, 0.82, size=14)

# カスケードフロー
cascade = [
    (C_INSLT, C_INS, "① 種類を選択", "特殊寝台 / 歩行器\n/車いす など13種"),
    (C_INSLT, C_INS, "② メーカーを選択", "種類で絞り込まれた\nメーカー一覧から選択"),
    (C_INSLT, C_INS, "③ 商品名を選択", "メーカーで絞り込まれた\n商品一覧から選択"),
    (C_ACCENTLT, C_ACCENT, "④ 自動入力", "タイスコード・単位数\nが自動でセット"),
]
x = 0.7
for bg, edge, step, desc in cascade:
    add_rect(s, x, 3.22, 2.78, 1.45, fill=bg, line=edge, lw=1)
    add_rect(s, x, 3.22, 0.09, 1.45, fill=edge)
    add_text(s, step, x + 0.22, 3.32, 2.45, 0.4, size=13, bold=True, color=C_INK)
    add_text(s, desc, x + 0.22, 3.75, 2.45, 0.78, size=12, color=C_MUTED)
    if x < 10:
        cx = x + 2.78
        add_rect(s, cx, 3.72, 0.28, 0.38, fill=C_LINE)
        add_text(s, "→", cx, 3.72, 0.28, 0.38, size=12, bold=True, color=C_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += 3.06

add_table(s, [
    ["フィールド", "入力内容", "必須"],
    ["種類（カテゴリ）", "特殊寝台 / 車いす / 歩行器 など13種類", "○"],
    ["メーカー", "マスターから選択（または手入力）", ""],
    ["商品名", "マスターから選択（または手入力）", "○"],
    ["タイスコード", "商品選択で自動入力（手入力も可）", ""],
    ["単位数", "商品選択で自動入力（手入力も可）", ""],
    ["卸会社", "仕入先の卸会社名", ""],
    ["利用開始日・終了日", "レンタル期間（カイポケCSV取込で自動設定）", ""],
    ["物件属性", "自社物件 / リース物件（Step 2で選択済み）", ""],
    ["カイポケ登録", "未登録 / 登録済（登録状況を記録）", ""],
], 0.7, 4.82, 11.9, 2.45, [2.6, 7.1, 1.2], header_color=C_INS)

# ════════════════════════════════════════════════
# 6. 自費レンタルの入力フォーム
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "自費レンタルの入力フォーム", eyebrow="SELF-PAY RENTAL FORM", bar=C_SELF)

add_rect(s, 0.7, 2.0, 11.9, 0.88, fill=C_SELFLT)
add_rect(s, 0.7, 2.0, 0.09, 0.88, fill=C_SELF)
add_lines(s, [
    ("自費レンタルは介護保険が使えない用具・期間のレンタルです。", False, C_INK),
    ("商品名・単価・数量・利用期間を入力します。税込金額は自動計算されます。", False, C_INK),
], 0.95, 2.08, 11.2, 0.7, size=14)

add_table(s, [
    ["フィールド", "入力内容", "備考"],
    ["商品名", "自費レンタルの品名", "自由入力"],
    ["卸会社", "仕入先の卸会社名", ""],
    ["単価", "月額レンタル料（税抜き）", "税込金額に自動換算"],
    ["数量", "レンタルする個数", "通常は1"],
    ["税区分", "非課税 / 10% / 軽8% / 税込", "選択式"],
    ["税込金額", "単価 × 数量 + 消費税", "自動計算（参考表示）"],
    ["利用開始日", "レンタル開始日", "月次売上集計に使用"],
    ["利用終了日", "レンタル終了日（空欄=継続中）", "空欄=当月の売上に含む"],
    ["物件属性", "自社物件 / リース物件（Step 2で選択済み）", ""],
    ["取引内容", "社内間取引 / ー", ""],
], 0.7, 3.1, 11.9, 3.55, [2.5, 6.3, 3.1], header_color=C_SELF)

callout(s, 0.7, 6.82, 11.9, 0.42,
        "🔒 売上確定済みの月の機器は編集・削除できません",
        "月次売上処理で「確定」した後は変更不可になります。修正が必要な場合は担当者に確定解除を依頼してください。",
        warn=True)

# ════════════════════════════════════════════════
# 7. 販売の入力フォーム
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "販売の入力フォーム", eyebrow="SALES FORM", bar=C_SALE)

add_rect(s, 0.7, 2.0, 11.9, 0.88, fill=C_SALELT)
add_rect(s, 0.7, 2.0, 0.09, 0.88, fill=C_SALE)
add_lines(s, [
    ("福祉用具の購入販売を記録します。納品日が月次売上処理の集計基準日になります。", False, C_INK),
    ("申請あり（日常生活給付等）の場合は、支払方法・申請市町村も記録してください。", False, C_INK),
], 0.95, 2.08, 11.2, 0.7, size=14)

add_table(s, [
    ["フィールド", "入力内容", "備考"],
    ["商品名（請求費目）", "販売品の品名", "必須"],
    ["受注日", "注文を受けた日", ""],
    ["納品日", "商品を届けた日", "必須（月次売上の集計基準）"],
    ["営業担当", "担当者名", ""],
    ["単価（税抜）", "商品の税抜単価", ""],
    ["数量", "販売個数", ""],
    ["税区分", "非課税 / 10% / 軽8% / 税込", "選択式"],
    ["送料（税抜）", "配送料金", "送料消費税は自動計算"],
    ["利用者自己負担割合", "全額負担 / 1〜3割 / 日常生活給付など", "選択式"],
    ["支払い方法", "口座引き落とし / 現金 / 受領委任払い など", "選択式"],
    ["申請あり", "チェックで申請状況を管理", "申請市町村も入力可"],
], 0.7, 3.1, 11.9, 3.9, [3.2, 5.8, 2.9], header_color=C_SALE)

# ════════════════════════════════════════════════
# 8. 機器の編集と削除
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "機器を編集する・削除する", eyebrow="EDIT / DELETE")

# 編集
add_rect(s, 0.7, 2.0, 5.75, 2.75, fill=C_INSLT, line=C_INS, lw=1)
add_rect(s, 0.7, 2.0, 0.09, 2.75, fill=C_INS)
add_text(s, "行クリックで編集する", 0.95, 2.12, 5.2, 0.48, size=18, bold=True, color=C_INS)
add_lines(s, [
    ("編集モード中に機器の行をクリックすると、", False, C_INK),
    ("詳細入力フォームが開きます。", False, C_INK),
    ("内容を修正して「保存」を押してください。", False, C_INK),
    ("（読み取りモード中はクリックしても開きません）", False, C_MUTED),
], 0.95, 2.65, 5.2, 1.9, size=14)

# 削除
add_rect(s, 6.95, 2.0, 5.65, 2.75, fill=C_WARNLT, line=C_WARN, lw=1)
add_rect(s, 6.95, 2.0, 0.09, 2.75, fill=C_WARN)
add_text(s, "削除ボタンで削除する", 7.2, 2.12, 5.1, 0.48, size=18, bold=True, color=C_WARN)
add_lines(s, [
    ("編集モード中、各行の右端にある", False, C_INK),
    ("ゴミ箱アイコン（赤）を押すと削除されます。", False, C_INK),
    ("確認なしに削除されるのでご注意ください。", True, C_WARN),
    ("（介護保険レンタルはカイポケCSVで復元可能）", False, C_MUTED),
], 7.2, 2.65, 5.1, 1.9, size=14)

# ロックの説明
add_rect(s, 0.7, 4.97, 11.9, 1.65, fill=RGBColor(0xF3, 0xF4, 0xF6), line=C_LINE, lw=1)
add_rect(s, 0.7, 4.97, 0.09, 1.65, fill=C_MUTED)
add_text(s, "🔒  確定済み機器は編集・削除できません", 0.95, 5.07, 11.2, 0.42, size=15, bold=True, color=C_INK)
add_lines(s, [
    ("月次売上処理で「売上確定」をした月に含まれる自費レンタル・販売機器は、", False, C_INK),
    ("一覧の商品名に 🔒 マークが付き、編集・削除できなくなります。", False, C_INK),
    ("修正が必要な場合は「売上・請求突合」ページで確定を解除してから操作してください。", True, C_WARN),
], 0.95, 5.55, 11.2, 0.95, size=13)

callout(s, 0.7, 6.77, 11.9, 0.4,
        "介護保険レンタルのロックについて",
        "介護保険レンタルもカイポケCSV確定後はロックされます。変更はカイポケ側で行いCSVを再インポートしてください。",
        ins=True)

# ════════════════════════════════════════════════
# 9. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="FAQ")

qas = [
    ("介護保険レンタルの機器を手動で追加してよいですか？",
     "追加はできますが、次回カイポケCSVインポート時に上書きされます。介護保険レンタルはカイポケCSVで管理するのが基本です。"),
    ("タイスコードが自動入力されません",
     "商品マスターに登録されていない機器は自動入力されません。「種類→メーカー→商品名」の順に選択し直すか、手入力してください。"),
    ("削除してしまいました。元に戻せますか？",
     "介護保険レンタルはカイポケCSVを再インポートすると復元できます。自費レンタル・販売は復元できません。内容を覚えている間に再入力してください。"),
    ("🔒マークが表示されて編集できません",
     "月次売上が確定済みのため編集不可です。売上・請求突合ページで当該月の確定を解除してから編集してください。"),
    ("自費レンタルと介護保険レンタルの両方に同じ機器があります",
     "利用区分の変更があった場合に両方に残ることがあります。不要な方の機器を削除して整理してください。"),
]
y = 2.0
rh = 0.96
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_ACCENT)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_ACCENT)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.38, size=14, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.44, 10.9, 0.44, size=12, color=C_MUTED)
    y += rh + 0.08

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "福祉用具選定タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
