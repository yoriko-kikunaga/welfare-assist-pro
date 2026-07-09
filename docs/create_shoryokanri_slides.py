"""
WelfareAssist Pro - 書類管理タブ 操作マニュアル（スライド生成）
出力: docs/書類管理タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_shoryokanri_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（他タブマニュアルと共通）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 書類管理固有色
C_SIGN     = RGBColor(0x7C, 0x3A, 0xED)   # 紫: サイン取得
C_SIGNLT   = RGBColor(0xED, 0xE9, 0xFE)
C_SIGNED   = RGBColor(0x16, 0xA3, 0x4A)   # 緑: 署名済み
C_SIGNEDLT = RGBColor(0xDC, 0xFC, 0xE7)
C_PLAN     = RGBColor(0x1D, 0x4E, 0xD8)   # 青: 計画書バッジ
C_PLANLT   = RGBColor(0xDB, 0xEA, 0xFE)
C_DELETE   = RGBColor(0xDC, 0x26, 0x26)   # 赤: 削除

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro ｜ 書類管理タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def callout(slide, l, t, w, h, title, body, warn=False, sign=False):
    if sign:
        bg, edge = C_SIGNLT, C_SIGN
    elif warn:
        bg, edge = C_WARNLT, C_WARN
    else:
        bg, edge = C_ACCENTLT, C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def doc_row(slide, y, badge_text, badge_bg, badge_fg, fname, meta,
            show_sign=False, show_signed_badge=False, rh=0.68):
    """1行の書類アイテムを描画"""
    add_rect(slide, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    # PDFアイコン
    add_rect(slide, 0.82, y + 0.1, 0.48, rh - 0.2, fill=RGBColor(0xFF,0xEE,0xEE))
    add_text(slide, "📄", 0.84, y + 0.12, 0.44, rh - 0.24, size=16, anchor=MSO_ANCHOR.MIDDLE)
    # 書類種別バッジ
    add_rect(slide, 1.44, y + 0.14, 0.72, 0.28, fill=badge_bg)
    add_text(slide, badge_text, 1.46, y + 0.14, 0.68, 0.28,
             size=10, bold=True, color=badge_fg, anchor=MSO_ANCHOR.MIDDLE)
    bx = 2.2
    if show_signed_badge:
        add_rect(slide, bx, y + 0.14, 0.88, 0.28, fill=C_SIGNEDLT)
        add_text(slide, "✓ 署名済み", bx + 0.02, y + 0.14, 0.84, 0.28,
                 size=10, bold=True, color=C_SIGNED, anchor=MSO_ANCHOR.MIDDLE)
    # ファイル名・メタ
    add_text(slide, fname, 1.44, y + 0.13, 6.2, 0.28, size=12, bold=True, color=C_INK)
    add_text(slide, meta, 1.44, y + 0.41, 6.2, 0.22, size=10, color=C_MUTED)
    # ボタン
    bx2 = 9.4
    if show_sign:
        add_rect(slide, bx2, y + 0.12, 1.1, 0.44, fill=C_SIGNLT, line=RGBColor(0xC4,0xB5,0xFD), lw=0.5)
        add_text(slide, "✍ サイン取得", bx2 + 0.04, y + 0.12, 1.02, 0.44,
                 size=10, bold=True, color=C_SIGN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bx2 += 1.2
    add_rect(slide, bx2, y + 0.12, 0.78, 0.44,
             fill=RGBColor(0xEF,0xF6,0xFF), line=RGBColor(0xBF,0xDB,0xFE), lw=0.5)
    add_text(slide, "↗ 開く", bx2 + 0.04, y + 0.12, 0.7, 0.44,
             size=10, bold=True, color=C_PLAN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx2 += 0.9
    add_rect(slide, bx2, y + 0.12, 0.72, 0.44,
             fill=RGBColor(0xFE,0xE2,0xE2), line=RGBColor(0xFE,0xCA,0xCA), lw=0.5)
    add_text(slide, "🗑 削除", bx2 + 0.04, y + 0.12, 0.64, 0.44,
             size=10, bold=True, color=C_DELETE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5,
         size=16, bold=True, color=C_ACCENT)
add_text(s, "「書類管理」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=46, bold=True, color=C_INK)
add_rect(s, 1.05, 3.95, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("計画書・モニタリングなどのPDF書類を保存・閲覧できる画面です。", False, RGBColor(0x5B, 0x6E, 0x68)),
    ("計画書への「サイン取得」機能で、タブレットから直接署名をもらうこともできます。", False, RGBColor(0x5B, 0x6E, 0x68)),
], 1.05, 4.25, 11, 1.2, size=17)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：利用者をクリック → 「書類管理」タブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. このタブでできること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")

feats = [
    (C_ACCENTLT, C_ACCENT, "📎", "書類をアップロード",
     "PDFファイルを種類（計画書・モニタリング・その他）を指定して保存。メモ入力も可。"),
    (RGBColor(0xEF,0xF6,0xFF), C_PLAN, "📂", "書類を開く・削除する",
     "保存済みの書類を新しいタブで確認。不要になったら削除できます（元に戻せません）。"),
    (C_SIGNLT, C_SIGN, "✍", "サイン取得（計画書のみ）",
     "計画書PDFを画面表示し、利用者やご家族にその場で手書きサインしてもらえます。"),
]
x = 0.7
for bg, edge, ic, title, body in feats:
    add_rect(s, x, 2.0, 3.8, 3.2, fill=bg, line=edge, lw=1)
    add_rect(s, x, 2.0, 0.1, 3.2, fill=edge)
    add_text(s, ic, x + 0.3, 2.2, 0.8, 0.7, size=28, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x + 0.28, 2.95, 3.3, 0.55, size=18, bold=True, color=C_INK)
    add_text(s, body, x + 0.28, 3.55, 3.35, 1.5, size=13, color=C_MUTED)
    x += 4.07

callout(s, 0.7, 5.55, 11.9, 0.9,
        "📝 書類の種類は「計画書」「モニタリング」「その他」の3種類",
        "サイン取得は計画書のみ使えます。モニタリング・その他はアップロードと閲覧・削除ができます。")

# ════════════════════════════════════════════════
# 3. 書類をアップロードする
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "書類をアップロードする", eyebrow="UPLOAD")

steps = [
    ("1", "「書類をアップロード」を押す", "画面右上の青いボタンを押すと、アップロードフォームが開きます。"),
    ("2", "種類を選んでPDFを選択", "書類の種類（計画書など）とメモ（任意）を入力し、「PDFファイルを選択」を押します。"),
    ("3", "完了・一覧に追加される", "アップロードが完了すると一覧に追加されます。クラウド保存なので全端末から見られます。"),
]
x = 0.7
for n, h, b in steps:
    add_rect(s, x, 2.1, 3.75, 2.6, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, n, x + 0.25, 2.25, 1.2, 1.0, size=44, bold=True, color=C_ACCENT)
    add_text(s, h, x + 0.25, 3.25, 3.3, 0.55, size=18, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.85, 3.3, 0.95, size=13, color=C_MUTED)
    x += 4.05

callout(s, 0.7, 5.05, 11.9, 1.05,
        "⚠ PDFファイルのみ・20MBまでです",
        "Word・Excel・画像は直接アップロードできません。「PDFとして保存」または「Export」でPDFに変換してからご利用ください。",
        warn=True)

# ════════════════════════════════════════════════
# 4. 書類一覧の見方
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "書類一覧の見方", eyebrow="LIST")

y = 1.85
doc_row(s, y, "計画書", C_PLANLT, C_PLAN,
        "福祉用具サービス計画書_2026-06_署名済み.pdf", "2026-07-09  ／  185 KB",
        show_signed_badge=True)
y += 0.78
doc_row(s, y, "計画書", C_PLANLT, C_PLAN,
        "福祉用具サービス計画書_2026-04.pdf", "2026-04-15  ／  142 KB  ／  4月分",
        show_sign=True)
y += 0.78
doc_row(s, y, "モニタリング", C_SIGNEDLT, C_SIGNED,
        "モニタリング報告書_2026-05.pdf", "2026-06-02  ／  98 KB")

add_table(s, [
    ["表示", "意味"],
    ["計画書・モニタリング・その他バッジ", "書類の種類"],
    ["✓ 署名済み（緑バッジ）", "サイン取得済みの書類"],
    ["✍ サイン取得（紫ボタン）", "計画書で署名版未作成の場合のみ表示"],
    ["↗ 開く（青ボタン）", "PDFを新しいタブで表示"],
    ["🗑 削除（赤ボタン）", "書類を削除（元に戻せません）"],
], 0.7, 5.6, 11.9, 1.75, [3.8, 8.1], size=13)

# ════════════════════════════════════════════════
# 5. 書類を開く・削除する
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "書類を開く・削除する", eyebrow="OPEN / DELETE")

# 開く
add_rect(s, 0.7, 2.0, 5.75, 2.4, fill=RGBColor(0xEF,0xF6,0xFF), line=RGBColor(0xBF,0xDB,0xFE), lw=1)
add_rect(s, 0.7, 2.0, 0.1, 2.4, fill=C_PLAN)
add_text(s, "↗  書類を開く", 0.95, 2.15, 5.2, 0.5, size=20, bold=True, color=C_PLAN)
add_lines(s, [
    ("「開く」ボタンを押すと、クラウドからPDFを取得して", False, C_INK),
    ("新しいタブで表示します。", False, C_INK),
    ("スマートフォン・タブレットでも同様に開けます。", False, C_MUTED),
], 0.95, 2.75, 5.2, 1.5, size=14)

# 削除
add_rect(s, 6.95, 2.0, 5.65, 2.4, fill=RGBColor(0xFE,0xE2,0xE2), line=RGBColor(0xFE,0xCA,0xCA), lw=1)
add_rect(s, 6.95, 2.0, 0.1, 2.4, fill=C_DELETE)
add_text(s, "🗑  削除する", 7.2, 2.15, 5.0, 0.5, size=20, bold=True, color=C_DELETE)
add_lines(s, [
    ("「削除」ボタンを押すと確認ダイアログが表示されます。", False, C_INK),
    ("「OK」を押すとクラウドから完全に削除されます。", False, C_INK),
    ("削除は元に戻せません。", True, C_DELETE),
], 7.2, 2.75, 5.1, 1.5, size=14)

callout(s, 0.7, 4.75, 11.9, 0.95,
        "💡 署名済みファイルと元の計画書は独立して管理されます",
        "署名済みファイルを削除しても元の計画書は残ります。逆に元の計画書を削除しても署名済みファイルは残ります。不要な方だけ削除してください。")

callout(s, 0.7, 5.9, 11.9, 0.75,
        "⚠ 誤って削除した場合は",
        "元のPDFファイルを再アップロードしてください（サイン取得が必要であれば再度行えます）。",
        warn=True)

# ════════════════════════════════════════════════
# 6. サイン取得とは
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "サイン取得機能とは（計画書のみ）", eyebrow="SIGN", bar=C_SIGN)

add_rect(s, 0.7, 2.0, 11.9, 1.35, fill=C_SIGNLT)
add_rect(s, 0.7, 2.0, 0.1, 1.35, fill=C_SIGN)
add_lines(s, [
    ("計画書PDFを画面に表示し、最終ページ右下の「利用者署名」欄に手書きでサインしてもらう機能です。", False, C_INK),
    ("タブレットやスマートフォンの画面上でそのまま指でサインできます。", False, C_INK),
    ("署名後は元のPDFとは別に「_署名済み.pdf」として保存されます（元ファイルは消えません）。", True, C_SIGN),
], 0.95, 2.1, 11.3, 1.15, size=15)

# 3つの特徴
feats3 = [
    (C_SIGNLT, C_SIGN, "その場でサイン",
     "訪問時にタブレットを渡すだけ。別途書類を持参・郵送する手間が省けます。"),
    (C_ACCENTLT, C_ACCENT, "元ファイルは保持",
     "署名済みPDFは新規作成されます。元の計画書は変わらず残ります。"),
    (C_SIGNEDLT, C_SIGNED, "自動でボタン管理",
     "署名済みファイルがあると元の計画書の「サイン取得」ボタンが自動的に消えます。"),
]
x = 0.7
for bg, edge, title, body in feats3:
    add_rect(s, x, 3.6, 3.8, 2.0, fill=bg, line=edge, lw=1)
    add_rect(s, x, 3.6, 0.08, 2.0, fill=edge)
    add_text(s, title, x + 0.25, 3.78, 3.3, 0.55, size=16, bold=True, color=C_INK)
    add_text(s, body, x + 0.25, 4.38, 3.35, 1.05, size=13, color=C_MUTED)
    x += 4.07

callout(s, 0.7, 5.85, 11.9, 0.8,
        "⚠ 署名の位置はA4計画書の最終ページ右下に固定されています",
        "計画書の様式が変わった場合や位置がずれる場合は担当者にご相談ください。",
        warn=True)

# ════════════════════════════════════════════════
# 7. サイン取得 操作手順
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "サイン取得 — 操作手順", eyebrow="SIGN STEPS", bar=C_SIGN)

steps4 = [
    ("1", "「サイン取得」を押す",
     "計画書（未署名・署名版未作成）の行にある紫のボタンを押します。PDFの読み込みが始まります。"),
    ("2", "PDFを確認する",
     "全ページが画面に表示されます。最終ページの右下に赤い枠で署名位置が示されます。"),
    ("3", "署名キャンバスに書く",
     "画面下部のキャンバスに指またはペンでサインします。「署名をクリア」でやり直しができます。"),
    ("4", "「署名して保存」を押す",
     "保存完了すると「_署名済み.pdf」として一覧に追加されます。元の計画書のボタンも自動的に消えます。"),
]
x = 0.7
for n, h, b in steps4:
    add_rect(s, x, 2.1, 2.85, 3.0, fill=C_WHITE, line=C_SIGNLT, lw=1.5)
    add_rect(s, x, 2.1, 0.1, 3.0, fill=C_SIGN)
    add_rect(s, x + 0.25, 2.22, 0.55, 0.55, fill=C_SIGN)
    add_text(s, n, x + 0.25, 2.22, 0.55, 0.55, size=20, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, h, x + 0.25, 2.88, 2.4, 0.6, size=15, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.53, 2.4, 1.45, size=12, color=C_MUTED)
    x += 3.07

callout(s, 0.7, 5.4, 11.9, 0.75,
        "💡 タブレット・スマートフォンで使うと便利です",
        "利用者訪問時にタブレットを差し出してそのまま指でサインしてもらえます。書いた内容はPDFに自動で埋め込まれます。",
        sign=True)

callout(s, 0.7, 6.35, 11.9, 0.75,
        "✍ サイン取得ボタンが表示されない場合",
        "「計画書」種別で、かつ署名済みファイルがまだ存在しない書類にのみ表示されます（他の種別・署名済みには表示されません）。",)

# ════════════════════════════════════════════════
# 8. 署名後の確認・やり直し
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "署名後の確認とやり直し", eyebrow="AFTER SIGN", bar=C_SIGN)

add_text(s, "署名が完了すると、書類一覧がこのように変わります。", 0.7, 2.0, 11.9, 0.45, size=15, color=C_MUTED)
y = 2.55
doc_row(s, y, "計画書", C_PLANLT, C_PLAN,
        "福祉用具サービス計画書_2026-06_署名済み.pdf", "2026-07-09  ／  185 KB",
        show_signed_badge=True)
y += 0.78
doc_row(s, y, "計画書", C_PLANLT, C_PLAN,
        "福祉用具サービス計画書_2026-06.pdf", "2026-07-01  ／  142 KB  ／  6月分")
# ← ボタンなし（署名版あるため自動非表示）

add_lines(s, [
    ("① 「_署名済み.pdf」が一覧の先頭（最新順）に追加される。", False, C_INK),
    ("② 元の計画書（2行目）には「サイン取得」ボタンが表示されなくなる。", False, C_INK),
    ("③ 署名済みファイルには「✓ 署名済み」バッジが表示される。", False, C_INK),
], 0.7, 4.28, 11.9, 0.9, size=14)

callout(s, 0.7, 5.35, 11.9, 0.85,
        "署名をやり直したい場合は",
        "署名済みファイルを削除してください。元の計画書に「サイン取得」ボタンが復活し、再度サインを取得できます。",
        sign=True)

callout(s, 0.7, 6.4, 11.9, 0.7,
        "⚠ 元の計画書を誤って削除しても署名済みファイルは残ります",
        "両ファイルは独立して管理されています。必要な方だけ削除してください。",
        warn=True)

# ════════════════════════════════════════════════
# 9. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="FAQ")

qas = [
    ("PDFしかアップロードできないのはなぜ？",
     "閲覧・印刷での再現性を保つためです。Word・ExcelはPDFとして保存してからご利用ください。"),
    ("「サイン取得」ボタンが表示されません",
     "書類の種類が「計画書」で、かつ署名済みファイルがまだ存在しない場合にのみ表示されます。"),
    ("書類は何件でも保存できますか？",
     "件数の上限はありませんが、1ファイル20MB以内です。古い書類は適宜削除してください。"),
    ("他の担当者がアップロードした書類は見られますか？",
     "はい。書類はクラウドに保存されており、ログイン済みのすべてのスタッフが確認できます。"),
    ("書類を誤って削除してしまいました",
     "元のPDFファイルを再アップロードしてください。"),
]
y = 2.0
rh = 0.95
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_ACCENT)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_ACCENT)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.38, size=14, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.44, 10.9, 0.44, size=12, color=C_MUTED)
    y += rh + 0.08

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "書類管理タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
