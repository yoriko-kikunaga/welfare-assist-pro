"""
WelfareAssist Pro - 売上・仕入突合タブ 操作マニュアル（スライド生成）
出力: docs/売上・仕入突合タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_reconciliation_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット ──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
C_INS      = RGBColor(0x25, 0x63, 0xEB)   # 青: 介護保険レンタル
C_INSLT    = RGBColor(0xDB, 0xEA, 0xFE)
C_SELF     = RGBColor(0x0D, 0x94, 0x88)   # ティール: 自費レンタル
C_SELFLT   = RGBColor(0xCC, 0xFB, 0xF1)
C_SALE     = RGBColor(0x7C, 0x3A, 0xED)   # 紫: 販売
C_SALELT   = RGBColor(0xED, 0xE9, 0xFE)
C_AMBER    = RGBColor(0xB4, 0x53, 0x09)
C_AMBERLT  = RGBColor(0xFE, 0xF3, 0xC7)
C_GREEN    = RGBColor(0x16, 0x65, 0x34)
C_GREENLT  = RGBColor(0xDC, 0xFC, 0xE7)
C_DELETE   = RGBColor(0xDC, 0x26, 0x26)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=27, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.58, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro ｜ 売上・仕入突合タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=12.5):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def callout(slide, l, t, w, h, title, body, warn=False):
    bg, edge = (C_WARNLT, C_WARN) if warn else (C_ACCENTLT, C_ACCENT)
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def step_cards(slide, steps, y=2.1, h=2.4):
    n = len(steps)
    gap = 0.28
    total_w = 11.9
    w = (total_w - gap * (n - 1)) / n
    x = 0.7
    for n_label, head, body in steps:
        add_rect(slide, x, y, w, h, fill=C_WHITE, line=C_LINE, lw=1)
        add_text(slide, n_label, x + 0.2, y + 0.1, 1.0, 0.75, size=32, bold=True, color=C_ACCENT)
        add_text(slide, head, x + 0.2, y + 0.85, w - 0.4, 0.55, size=13.5, bold=True, color=C_INK)
        add_text(slide, body, x + 0.2, y + 1.35, w - 0.4, h - 1.45, size=11, color=C_MUTED)
        x += w + gap


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 1.85, 11, 0.5, size=16, bold=True, color=C_ACCENT)
add_text(s, "「売上・仕入突合」タブの使い方", 1.0, 2.45, 11.3, 1.3, size=42, bold=True, color=C_INK)
add_rect(s, 1.05, 3.75, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("卸会社からの請求書（仕入）と社内の売上データを、", False, C_MUTED),
    ("利用者・品目単位で照合し粗利を確認する画面です。", False, C_MUTED),
    ("7社の請求書アップロード、AIによる自動読み取り、", False, C_MUTED),
    ("利用者ごとの品目突合まで一通り行えます。", False, C_MUTED),
], 1.05, 4.05, 11, 1.8, size=16)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：左サイドバー「売上・仕入突合」（青緑ボタン）",
         1.05, 6.5, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. できること（3ステップ）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")
add_text(s, "流れは「アップロード → 突合 → 確定」の3ステップです。", 0.7, 1.85, 11.9, 0.4, size=14, color=C_MUTED)
step_cards(s, [
    ("1", "請求書をアップロード", "卸会社ごとにPDFまたはCSVをアップロード。PDFはAIが自動でテキストを読み取ります。"),
    ("2", "突合・紐づけ", "自動で利用者名がマッチング。未マッチは画面上で手動紐づけ。粗利・差額を確認します。"),
    ("3", "確定", "会社ごとに仕入確定→全社・全売上がそろったら月次確定。確定後はデータが固定されます。"),
], y=2.4, h=2.6)

# ════════════════════════════════════════════════
# 3. 画面の3メインタブ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "画面の3つのメインタブ", eyebrow="MAIN TABS")
add_rect(s, 0.7, 2.0, 2.9, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "売上一覧 (1822)", 0.7, 2.0, 2.9, 0.55, size=13, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 3.75, 2.0, 3.3, 0.55, fill=C_ACCENT)
add_text(s, "請求書アップロード (5)", 3.75, 2.0, 3.3, 0.55, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 7.2, 2.0, 2.6, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "突合結果", 7.2, 2.0, 2.6, 0.55, size=13, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_table(s, [
    ["メインタブ", "できること"],
    ["売上一覧", "選んだ月・事業所の当月売上（介護保険レンタル・自費レンタル・販売）を一覧確認するだけの画面"],
    ["請求書アップロード", "卸会社ごとに請求書PDF/CSVをアップロード・確認・「突合実行」ボタン"],
    ["突合結果", "紐づけの修正、利用者別突合セクション、CSV出力、確定作業。「突合実行」を押すまでは選べない（グレーアウト）"],
], 0.7, 2.9, 11.9, 1.9, [2.8, 9.1], size=13)

callout(s, 0.7, 5.0, 11.9, 1.5,
        "⚠ 「突合結果」タブがグレーで押せないときは異常ではありません",
        "まだ「突合実行」ボタンを押していないだけです。月度や事業所を切り替えるたびにこの結果はリセットされます。その都度「請求書アップロード」タブ下部の緑色「突合実行」ボタンを押してください。",
        warn=True)

# ════════════════════════════════════════════════
# 4. 対応卸会社7社
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "対応している卸会社（7社）", eyebrow="COMPANIES")
companies = [
    ("日建リース工業株式会社", "PDF"), ("野口株式会社", "PDF"),
    ("株式会社ニシケン", "CSV/PDF"), ("パラマウントケアサービス", "CSV/PDF"),
    ("日本ケアサプライ", "CSV/PDF"), ("株式会社キシヤ", "PDF"),
    ("その他", "PDF/画像"),
]
x, y = 0.7, 2.0
for i, (name, tag) in enumerate(companies):
    col = i % 2
    row = i // 2
    cx = 0.7 + col * 6.0
    cy = 2.0 + row * 0.85
    add_rect(s, cx, cy, 5.8, 0.65, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, name, cx + 0.2, cy, 3.6, 0.65, size=13, bold=True, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    tagcol = C_GREENLT if "CSV" in tag else C_INSLT
    tagfg = C_GREEN if "CSV" in tag else C_INS
    add_rect(s, cx + 4.0, cy + 0.13, 1.6, 0.4, fill=tagcol)
    add_text(s, tag, cx + 4.0, cy + 0.13, 1.6, 0.4, size=11, bold=True, color=tagfg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

callout(s, 0.7, 6.05, 11.9, 1.0,
        "💡 3社はCSV取込にも対応",
        "ニシケン・パラマウント・日本ケアサプライはCSVでの取込みが可能です（AIを使わず高速・正確）。可能な場合はCSV取込みをおすすめします。")

# ════════════════════════════════════════════════
# 5. 請求書アップロード
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "請求書をアップロードする", eyebrow="UPLOAD")
add_text(s, "「請求書アップロード」タブで、卸会社ごとのカードにファイルをドラッグ＆ドロップまたは選択します。",
         0.7, 1.85, 11.9, 0.5, size=13.5, color=C_MUTED)

add_rect(s, 0.7, 2.5, 5.75, 1.9, fill=C_INSLT)
add_rect(s, 0.7, 2.5, 0.08, 1.9, fill=C_INS)
add_lines(s, [
    ("📄 PDF・画像の場合", True, C_INS),
    ("AIが自動でテキストを読み取ります（約30〜60秒）。", False, C_INK),
    ("同じ会社のPDFを複数回に分けてアップロードすると自動マージ。", False, C_INK),
    ("金額記載のない請求書（キシヤ等）は仕入金額0円扱い。", False, C_INK),
], 0.95, 2.65, 5.3, 1.65, size=12.5)

add_rect(s, 6.85, 2.5, 5.75, 1.9, fill=C_GREENLT)
add_rect(s, 6.85, 2.5, 0.08, 1.9, fill=C_GREEN)
add_lines(s, [
    ("📊 CSVの場合", True, C_GREEN),
    ("ニシケン・パラマウント・日本ケアサプライの3社対応。", False, C_INK),
    ("AIを使わず自動でパースされ、速く正確です。", False, C_INK),
], 7.1, 2.65, 5.3, 1.65, size=12.5)

callout(s, 0.7, 4.6, 11.9, 1.1,
        "✅ アップロード後に金額の検証結果が表示されます",
        "「検証OK：請求書合計と一致」＝AI読み取りの合計と請求書の総額が一致（差額1,000円以内）。CSV取込は請求書合計自体がないため「CSV取込（請求書合計なし）」とグレー表示されます。")

callout(s, 0.7, 5.85, 11.9, 0.95,
        "🔒 確定済みの卸会社にはアップロードできません",
        "「確定済みのため編集不可」と出た場合は、先にその会社の仕入確定を解除してください。",
        warn=True)

# ════════════════════════════════════════════════
# 6. 突合実行
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "突合実行", eyebrow="RUN MATCHING")
add_text(s, "請求書アップロードタブの一番下にある緑色の大きなボタンです。", 0.7, 2.0, 11.9, 0.4, size=14, color=C_MUTED)
add_rect(s, 0.7, 2.7, 8.0, 0.75, fill=C_ACCENT)
add_text(s, "⇄  突合実行（売上1822件 × 請求書5社）", 0.7, 2.7, 8.0, 0.75, size=16, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
callout(s, 0.7, 3.85, 11.9, 1.2,
        "このボタンを押すと",
        "その時点の当月売上データとアップロード済み全社の請求書を照合し、「突合結果」タブに結果が表示されます。ボタンはアップロード済みの請求書が1社以上あれば押せます。")

# ════════════════════════════════════════════════
# 7. 突合結果タブの構成（重要仕様）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "突合結果タブの構成（重要な仕様）", eyebrow="RESULT TABS")
add_text(s, "「突合結果」タブの中には、さらに3つのサブタブがあります。", 0.7, 1.85, 11.9, 0.4, size=13.5, color=C_MUTED)

add_rect(s, 0.7, 2.35, 3.3, 0.55, fill=C_ACCENT)
add_text(s, "突合済み (2173)", 0.7, 2.35, 3.3, 0.55, size=12.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 4.15, 2.35, 3.0, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "売上のみ (261)", 4.15, 2.35, 3.0, 0.55, size=12.5, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 7.3, 2.35, 3.0, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "仕入のみ (243)", 7.3, 2.35, 3.0, 0.55, size=12.5, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, 0.7, 3.15, 11.9, 1.7, fill=C_AMBERLT)
add_rect(s, 0.7, 3.15, 0.08, 1.7, fill=C_AMBER)
add_lines(s, [
    ("📌 件数バッジに数字があっても「突合済み」「売上のみ」の表は空欄が正常です", True, C_AMBER),
    ("このアプリの売上は介護保険レンタル・自費レンタル・販売の3種類のみですが、", False, C_INK),
    ("これら3種類は下部の専用「利用者別突合セクション」で個別に管理される仕組みのため、", False, C_INK),
    ("この上部の表からは常に除外されます。結果として2タブは実質いつも空欄になります。", False, C_INK),
], 0.95, 3.3, 11.4, 1.5, size=13)

callout(s, 0.7, 5.05, 11.9, 0.9,
        "✅ 「仕入のみ」タブだけは実際に使います",
        "まだどの利用者にも紐づいていない、または当月の売上に対応が見つからない仕入データが表示されます。")

# ════════════════════════════════════════════════
# 8. インライン紐づけ編集
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "インライン紐づけ編集", eyebrow="INLINE EDIT")
add_text(s, "CSV出力→修正→再インポートの手間なく、画面上でそのまま紐づけを直せます。", 0.7, 1.85, 11.9, 0.4, size=13.5, color=C_MUTED)

add_rect(s, 0.7, 2.4, 11.9, 0.6, fill=C_INSLT)
add_text(s, "瓜生健次様 → 瓜生 健次 (9155)", 0.9, 2.4, 4.5, 0.6, size=12, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "¥0", 5.6, 2.4, 1.0, 0.6, size=12, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "日建リース工業株式会社", 6.8, 2.4, 2.8, 0.6, size=11, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 9.8, 2.53, 0.9, 0.35, fill=C_GREENLT)
add_text(s, "変更", 9.8, 2.53, 0.9, 0.35, size=10.5, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 10.8, 2.53, 1.0, 0.35, fill=C_AMBERLT)
add_text(s, "対象月", 10.8, 2.53, 1.0, 0.35, size=10.5, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_table(s, [
    ["タブ", "操作", "できること"],
    ["仕入のみ", "「紐づけ」／「変更」ボタン", "未紐づけの仕入を利用者に紐づける。紐づけ済みは「変更」（緑）で変更・解除"],
    ["突合済み", "鉛筆アイコン", "紐づけの変更・解除"],
    ["売上のみ", "「仕入紐づけ」ボタン", "売上に対応する仕入データを手動で選択"],
], 0.7, 3.25, 11.9, 1.9, [2.0, 3.4, 6.5], size=12.5)

callout(s, 0.7, 5.35, 11.9, 1.15,
        "🔒 確定済みの卸会社は編集不可",
        "その卸会社の請求書が「確定済み」の間はボタンが薄いグレーで押せません。先に「請求書アップロード」タブで「解除」を押してください。解除しても請求書の金額データは消えません。",
        warn=True)

# ════════════════════════════════════════════════
# 9. 対象月度タグ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "対象月度タグ（月をまたぐ遅れ請求）", eyebrow="TARGET MONTH", bar=C_AMBER)
add_text(s, "卸会社が前月以前の売上分をまとめて今月の請求書に含めてくることがあります。品目ごとに「本来の対象月度」を指定できる機能です。",
         0.7, 1.85, 11.9, 0.6, size=13, color=C_MUTED)

step_cards(s, [
    ("1", "「対象月」ボタンを押す", "「仕入のみ」タブの行にある黄色い「対象月」ボタンを押すと、月選択の入力欄が現れます。"),
    ("2", "本来の月度を選ぶ", "例：今月（6月）にアップロードされた請求書だが、実際は5月分の売上に対する請求なら「2026年05月」を選びます。"),
], y=2.6, h=1.9)

add_rect(s, 0.7, 4.7, 11.9, 0.95, fill=C_ACCENTLT)
add_rect(s, 0.7, 4.7, 0.08, 0.95, fill=C_ACCENT)
add_lines(s, [
    ("設定すると利用者名の横に「対象:2026年05月」バッジが表示されます。", True, C_ACCENT),
    ("「販売」「自費レンタル」の利用者別突合セクションでは、その月の売上データも合算して差額を再計算します。今月の請求書合計は変わりません。", False, C_INK),
], 0.95, 4.82, 11.4, 0.8, size=12.5)

add_rect(s, 0.7, 5.8, 11.9, 1.3, fill=C_WARNLT)
add_rect(s, 0.7, 5.8, 0.08, 1.3, fill=C_WARN)
add_lines(s, [
    ("⚠ 2つの制約：①対象月度側の画面を直接開いても自動では表示されません（アップロードした月の画面で処理した場合のみ正しく見える）", True, C_WARN),
    ("②介護保険レンタルの給付対象金額は毎月上書きされる単一値のため、過去月の正確な金額を遡って再現できません。", False, C_INK),
], 0.95, 5.92, 11.4, 1.15, size=12)

# ════════════════════════════════════════════════
# 10. 利用者別突合セクション：対象条件
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "利用者別突合セクション（1）対象条件", eyebrow="CLIENT SECTIONS")
add_text(s, "突合結果タブの下部に、3種類の専用セクションが表示されます。実務ではここが突合作業の中心です。",
         0.7, 1.85, 11.9, 0.5, size=13.5, color=C_MUTED)

add_table(s, [
    ["セクション", "色", "対象利用者"],
    ["介護保険レンタル", "青", "当月アクティブな介護保険レンタルがある利用者（＋対象月度タグ分）"],
    ["販売", "紫", "当月納品日の販売品目がある利用者（＋対象月度タグ分）"],
    ["自費レンタル", "ティール", "当月アクティブな自費レンタルがある利用者（＋対象月度タグ分）"],
], 0.7, 2.5, 11.9, 1.9, [3.0, 1.7, 7.2], size=13)

# 会社ブロックのモック
add_rect(s, 0.7, 4.6, 11.9, 0.85, fill=RGBColor(0xF9,0xFA,0xFB), line=C_LINE, lw=0.5)
add_text(s, "株式会社日本ケアサプライ", 0.9, 4.65, 4.0, 0.35, size=13, bold=True, color=C_INK)
add_rect(s, 5.0, 4.72, 1.6, 0.32, fill=RGBColor(0xFE,0xE2,0xE2))
add_text(s, "差額あり 8件", 5.0, 4.72, 1.6, 0.32, size=10.5, bold=True, color=C_DELETE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "8名 ／ 弊社計：¥232,400 ／ 卸計：¥65,665", 0.9, 5.05, 6.0, 0.3, size=11, color=C_MUTED)
add_rect(s, 10.3, 4.72, 1.5, 0.42, fill=C_SALE)
add_text(s, "確定する", 10.3, 4.72, 1.5, 0.42, size=11.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

callout(s, 0.7, 5.65, 11.9, 1.15,
        "🖱️ 会社名クリックで利用者一覧が展開",
        "差額がある行は赤背景、紐づけ保存済みは緑背景＋「紐づけ済」バッジ。行の「詳細」／「確認・編集」ボタンで品目突合モーダルが開きます。")

# ════════════════════════════════════════════════
# 11. 品目突合モーダル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "利用者別突合セクション（2）品目突合モーダル", eyebrow="ITEM MATCHING")

add_rect(s, 0.7, 1.95, 11.9, 2.6, fill=C_WHITE, line=C_LINE, lw=1)
add_text(s, "サービス付き高齢者向け住宅あおぞら東千石 ー 品目突合", 0.9, 2.05, 11.5, 0.4, size=13, bold=True, color=C_INK)
add_rect(s, 0.7, 2.5, 11.9, 0.4, fill=RGBColor(0xF3,0xF4,0xF6))
add_text(s, "販売合計：¥3,800　　卸請求合計：¥0　　差額：+3,800円", 0.9, 2.5, 11.5, 0.4, size=11.5, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 0.9, 3.1, 10.9, 0.85, fill=RGBColor(0xF9,0xFA,0xFB))
add_text(s, "優肌絆 EasyCutホワイト 25mm×7m", 1.05, 3.25, 4.0, 0.55, size=12, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "→", 5.3, 3.25, 0.4, 0.55, size=16, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 5.9, 3.25, 5.6, 0.55, fill=C_WHITE, line=C_GREENLT, lw=1)
add_text(s, "優肌絆 EasyCutホワイト MM3262　¥4,568", 6.05, 3.25, 5.4, 0.55, size=11.5, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)

callout(s, 0.7, 4.75, 11.9, 1.55,
        "🖱️ できること",
        "「卸品目を追加」（ドロップダウン選択）で品目追加、卸品目タグ右の「×」で解除。1つの弊社品目に複数の卸品目を紐づけ可能（1:N対応・附属品など）。「紐づけを保存」を押すと翌月以降も自動的に同じ紐づけが適用されます。")

# ════════════════════════════════════════════════
# 12. 自社ベッド表示の注意
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "自社ベッド（仕入不要）の表示", eyebrow="COMPANY-OWNED")
add_lines(s, [
    ("機器編集モーダルで「物件属性」を「自社物件」に設定した品目は、", False, C_INK),
    ("品目突合モーダルで紫色の背景＋「自社ベッド」バッジとなり、", False, C_INK),
    ("「仕入不要（自社ベッド）」と表示されます（卸品目追加ボタンは出ません）。", False, C_INK),
    ("この品目は卸請求合計の計算から自動的に除外されます。", True, C_SALE),
], 0.7, 1.95, 11.9, 1.7, size=14)

add_rect(s, 0.7, 3.8, 11.9, 1.3, fill=C_WARNLT)
add_rect(s, 0.7, 3.8, 0.08, 1.3, fill=C_WARN)
add_lines(s, [
    ("⚠ 「物件属性」はベッドに限らず、どの機器種別にも設定できます", True, C_WARN),
    ("ベッド以外の消耗品などに誤って設定すると、実際には仕入が発生している品目が", False, C_INK),
    ("「仕入不要」扱いになり、卸請求合計から漏れてしまいます。品目突合モーダルで", False, C_INK),
    ("「自社ベッド」バッジがついた品目が本当に自社所有物件か、時々見直してください。", False, C_INK),
], 0.95, 3.92, 11.4, 1.15, size=12.5)

callout(s, 0.7, 5.3, 11.9, 1.1,
        "会社単位の確定",
        "セクション内の全利用者を確認し終えたら「確定する」ボタン（紫）で会社単位を確定。確定後は「確定解除」で解除できます。")

# ════════════════════════════════════════════════
# 13. 卸品目未紐づけ一覧
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "卸品目未紐づけ一覧", eyebrow="UNMATCHED CHECK")
add_text(s, "利用者別突合セクションのさらに下部にある、月次確定前の最終チェック用セクションです。",
         0.7, 1.9, 11.9, 0.4, size=13.5, color=C_MUTED)
step_cards(s, [
    ("1", "「読み込む」ボタンを押す", "Firestoreから紐づけ情報を取得し、まだ弊社品目に紐づいていない卸請求品目を一覧表示します。"),
    ("2", "紐づけ漏れを確認", "種別・利用者名・施設名・弊社品目・卸品目・卸金額・卸会社が一覧に。CSV出力も可能です。"),
], y=2.5, h=1.9)
callout(s, 0.7, 4.6, 11.9, 1.0,
        "💡 一覧の更新は自動ではありません",
        "紐づけ作業を進めたあとは、もう一度「読み込む」を押して一覧を更新してください。")

# ════════════════════════════════════════════════
# 14. 確定の3段階フロー
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "確定の流れ（3段階）", eyebrow="CONFIRM FLOW")
steps_conf = [
    ("1", "売上確定（3種類）", "「月次売上処理」タブで介護保険レンタル・自費レンタル・販売をそれぞれ確定"),
    ("2", "仕入確定（卸会社ごと）", "この画面でアップロード済みの卸会社を1社ずつ確定"),
    ("3", "月次確定", "売上3種類＋仕入全社が確定済みのときだけ押せる、その月度の最終確定"),
]
y = 1.95
for n, h, b in steps_conf:
    add_rect(s, 0.7, y, 11.9, 1.05, fill=C_WHITE, line=C_LINE, lw=1)
    add_rect(s, 0.9, y + 0.18, 0.5, 0.5, fill=C_ACCENT)
    add_text(s, n, 0.9, y + 0.18, 0.5, 0.5, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, h, 1.6, y + 0.12, 9.5, 0.4, size=15, bold=True, color=C_INK)
    add_text(s, b, 1.6, y + 0.52, 9.9, 0.45, size=12, color=C_MUTED)
    y += 1.2
    if n != "3":
        add_text(s, "↓", 0.9, y - 0.15, 0.5, 0.3, size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

add_rect(s, 0.7, 5.65, 11.9, 1.35, fill=C_ACCENTLT)
add_rect(s, 0.7, 5.65, 0.08, 1.35, fill=C_ACCENT)
add_lines(s, [
    ("🔓 解除は逆順（月次→個別）です", True, C_ACCENT),
    ("月次確定を解除しないと個別の売上確定・仕入確定は解除できません。", False, C_INK),
    ("「解除」自体はデータや金額を変更しません。修正のために一時的に解除するのは安全です。", False, C_INK),
], 0.95, 5.78, 11.4, 1.2, size=12.5)

# ════════════════════════════════════════════════
# 15. CSV出力と再インポート
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "CSV出力と再インポート", eyebrow="CSV")
step_cards(s, [
    ("1", "CSV出力", "「突合結果」タブの「CSVエクスポート」ボタンで、統一ヘッダー（あおぞらID・利用者名・商品名・種別・売上金額・仕入金額・粗利・粗利率・卸会社）のCSVを保存。"),
    ("2", "修正して再インポート", "ExcelであおぞらIDを追記・修正し、「突合CSVインポート」で取り込むと紐づけが反映されます。あおぞらIDがある行は名前マッチングをスキップして直接紐づけ。"),
], y=2.3, h=2.2)

# ════════════════════════════════════════════════
# 16. 毎月のチェックリスト
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "毎月の作業チェックリスト", eyebrow="CHECKLIST")
items = [
    "「月次売上処理」タブで3種類の売上を確定した",
    "対象月の請求書を全社分アップロードし、検証OKになっているか確認した",
    "「突合実行」を押して突合結果タブを開けるようにした",
    "利用者別突合セクション（3種類）で差額のある利用者を確認・紐づけ修正した",
    "月をまたぐ遅れ請求があれば「対象月」でタグ付けした",
    "「卸品目未紐づけ一覧」を確認し、紐づけ漏れがないか見た",
    "卸会社を1社ずつ仕入確定し、全社・全売上がそろったら月次確定した",
]
y = 1.9
for it in items:
    add_rect(s, 0.7, y, 11.9, 0.58, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.9, y + 0.15, 0.28, 0.28, fill=None, line=C_ACCENT, lw=1.5)
    add_text(s, it, 1.35, y, 11.0, 0.58, size=12.5, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.67

# ════════════════════════════════════════════════
# 17. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問・トラブルシューティング", eyebrow="FAQ")
qas = [
    ("「突合結果」タブがグレーで押せません",
     "「請求書アップロード」タブに戻り、下部の緑色「突合実行」ボタンを押してください。"),
    ("「突合済み」「売上のみ」タブが空欄です",
     "正常です。3種類の売上は下の「利用者別突合セクション」で管理されるため、この2タブは常に空欄になります。"),
    ("紐づけを変更したいのにボタンが押せません",
     "その卸会社の請求書が「確定済み」です。「解除」を押してから再度お試しください。金額データは消えません。"),
    ("前月分の遅れ請求が混ざっていました",
     "「仕入のみ」タブでその行に「対象月」を設定してください。今月の請求合計は変わりません。"),
    ("消耗品が「自社ベッド」と表示されています",
     "機器編集で「物件属性」を誤って「自社物件」に設定した可能性があります。「ー」に戻してください。"),
]
y = 1.9
rh = 0.97
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_ACCENT)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_ACCENT)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.4, size=13, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.45, 10.9, 0.47, size=11.5, color=C_MUTED)
    y += rh + 0.06

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "売上・仕入突合タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
