"""
WelfareAssist Pro - 基本情報タブ 操作マニュアル（スライド生成）
出力: docs/基本情報タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_kihon_joho_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（アプリ／HTMLマニュアルと対応）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)  # 淡いセージ
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)  # 深い松葉色
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)  # ティール（識別色）
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)  # 注意の暖色
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 機能色（画面上の色と対応）
C_FACILITY = RGBColor(0x0E, 0x8A, 0x78)
C_CM       = RGBColor(0x25, 0x63, 0xEB)
C_RECEIPT  = RGBColor(0xE1, 0x1D, 0x67)
C_WELFARE  = RGBColor(0x16, 0xA3, 0x4A)
C_KAIPOKE  = RGBColor(0xB4, 0x53, 0x09)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, color=C_INK, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)           # 左帯
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro ｜ 基本情報タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def callout(slide, l, t, w, h, title, body, warn=False):
    bg = C_WARNLT if warn else C_ACCENTLT
    edge = C_WARN if warn else C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5,
         size=16, bold=True, color=C_ACCENT)
add_text(s, "「基本情報」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=46, bold=True, color=C_INK)
add_rect(s, 1.05, 3.95, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("利用者ひとりひとりの基本データを登録・確認する画面です。", False, C_MUTED),
    ("はじめての方でも迷わず編集・保存できるよう、やさしくまとめました。", False, C_MUTED),
], 1.05, 4.25, 11, 1.2, size=17)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：利用者をクリック → 一番左のタブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. このタブでできること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")
add_rect(s, 0.7, 2.0, 11.9, 1.5, fill=C_ACCENTLT)
add_rect(s, 0.7, 2.0, 0.08, 1.5, fill=C_ACCENT)
add_lines(s, [
    ("このタブは、利用者の「カルテの表紙」のような画面です。", True, C_INK),
    ("請求・レセプトチェック・売上集計など、アプリのほかの機能はすべて", False, C_INK),
    ("ここの情報をもとに動きます。正しく入れておくほど、後の作業がラクになります。", False, C_INK),
], 1.0, 2.2, 11.3, 1.3, size=16)
add_text(s, "この画面で扱う主な情報", 0.7, 3.9, 11, 0.45, size=16, bold=True, color=C_ACCENT)
add_table(s, [
    ["分類", "おもな項目"],
    ["本人・住まい", "氏名・フリガナ・生年月日・性別・入居施設名・居室番号・在宅区分"],
    ["事業・福祉用具", "事業所・福祉用具利用者・レセプトチェック対象・請求区分"],
    ["介護保険", "要介護度・負担割合・被保険者証・負担割合証・支払い区分"],
    ["関係先・施設", "ケアマネ情報・施設契約情報（入退去）・キーパーソン・カイポケ登録"],
], 0.7, 4.35, 11.9, 2.3, [2.6, 9.3])

# ════════════════════════════════════════════════
# 3. まず3ステップ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "まず、この3ステップだけ", eyebrow="BASIC")
steps = [
    ("1", "「編集」を押す", "画面右上の「編集」ボタンを押すと、入力できる状態になります。"),
    ("2", "内容を入力する", "直したい欄を選んで入力・選択します。複数まとめてOK。"),
    ("3", "「保存する」を押す", "右上の「保存する」で確定。これで全員に反映されます。"),
]
x = 0.7
for n, h, b in steps:
    add_rect(s, x, 2.1, 3.75, 2.4, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, n, x + 0.25, 2.25, 1.2, 1.0, size=44, bold=True, color=C_ACCENT)
    add_text(s, h, x + 0.25, 3.25, 3.3, 0.5, size=18, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.75, 3.3, 0.9, size=13, color=C_MUTED)
    x += 4.05
callout(s, 0.7, 4.85, 11.9, 1.0,
        "⚠️ 自動保存はありません（2026-08〜）",
        "編集したら必ず右上の「保存する」を押してください。未保存の間はヘッダーに「未保存の変更があります」と表示され、"
        "保存せず他の利用者に移動・タブを閉じようとすると確認が出ます。保存内容は毎日の自動更新のあとも残ります。",
        warn=True)

# ════════════════════════════════════════════════
# 4. 画面の地図
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "画面の地図（上から下への並び）", eyebrow="MAP")
rows_map = [
    ("1", "あおぞらID・事業所", "利用者番号と担当事業所", C_ACCENT, True),
    ("2", "氏名・フリガナ・生年月日・性別", "本人のプロフィール", C_INK, False),
    ("3", "施設契約情報", "施設への入居・退去（Kintone自動連携）", C_FACILITY, False),
    ("4", "入居施設名・居室番号・在宅区分", "住まいの区分", C_INK, True),
    ("5", "福祉用具利用者／レセプト対象／請求区分", "福祉用具・請求にかかわる設定", C_WELFARE, True),
    ("6", "ケアマネージャー情報", "居宅介護支援事業所・担当CM", C_CM, True),
    ("7", "介護保険情報", "要介護度・負担割合・各種証", C_ACCENT, True),
    ("8", "支払い区分", "生保・非生保", C_INK, True),
    ("9", "キーパーソン", "ご家族など連絡先", C_INK, False),
    ("10", "カイポケ登録（基本情報）", "カイポケへの登録状況", C_KAIPOKE, False),
]
y = 1.85
rh = 0.5
for n, title, desc, col, clock in rows_map:
    add_rect(s, 0.7, y, 0.42, rh - 0.08, fill=col)
    add_text(s, n, 0.7, y - 0.02, 0.42, rh - 0.08, size=15, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, 1.25, y - 0.02, 6.5, rh, size=15, bold=True, color=C_INK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, desc, 7.6, y - 0.02, 4.0, rh, size=12, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    if clock:
        add_text(s, "🕐 履歴あり", 11.5, y - 0.02, 1.4, rh, size=11, bold=True,
                 color=C_ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    y += rh
add_text(s, "※「🕐」がある項目は、いつから変わったかの履歴を残せます（最後の章参照）。",
         1.25, y + 0.05, 11, 0.4, size=12, color=C_MUTED)

# ════════════════════════════════════════════════
# 5. 項目 1・2
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "項目くわしく ①②　基本プロフィール", eyebrow="FIELDS")
add_text(s, "① あおぞらID・事業所", 0.7, 1.85, 11, 0.45, size=16, bold=True, color=C_ACCENT)
add_table(s, [
    ["項目", "入れる内容", "選べる値／形式"],
    ["あおぞらID", "利用者を見分ける番号。基本は変更しません", "AZ-0001 などの形式"],
    ["事業所 🕐", "担当事業所。売上・請求はこの事業所で集計", "鹿児島（ACG）／福岡（Lichi）"],
], 0.7, 2.3, 11.9, 1.3, [2.4, 6.0, 3.5])
add_text(s, "② 氏名・フリガナ・生年月日・性別", 0.7, 3.95, 11, 0.45, size=16, bold=True, color=C_ACCENT)
add_table(s, [
    ["項目", "入れる内容", "選べる値／形式"],
    ["氏名", "利用者のお名前", "文字入力"],
    ["フリガナ", "検索・あいうえお順の並び替えに使用", "カナ入力"],
    ["生年月日", "カレンダーから選択", "日付"],
    ["性別", "選択", "男性／女性"],
], 0.7, 4.4, 11.9, 2.1, [2.4, 6.0, 3.5])

# ════════════════════════════════════════════════
# 6. 項目 3 施設契約情報
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "項目くわしく ③　施設契約情報", eyebrow="FIELDS", bar=C_FACILITY)
add_lines(s, [
    ("施設への「入居」「退去」の記録です。Kintone から自動で入ってくる情報で、", False, C_INK),
    ("入居と退去がペアで表示されます。施設に入っている利用者だけ表示されます。", False, C_INK),
], 0.7, 1.85, 11.9, 0.9, size=15)
add_table(s, [
    ["表示", "意味"],
    ["入居日", "施設に入った日（自動連携）"],
    ["退去日", "施設を出た日。入居中なら「退去情報なし（入居中）」と表示"],
    ["記録者・特記", "連携元の記録者名と補足メモ"],
], 0.7, 2.75, 11.9, 1.7, [2.6, 9.3], header_color=C_FACILITY)
callout(s, 0.7, 4.7, 11.9, 1.9,
        "⚠️ ここはレンタル契約とは別物です",
        "福祉用具レンタルの「新規・解約」は、別タブ（利用者新規・変更情報入力）の"
        "「レンタル契約情報」にあります。混同にご注意ください。\n"
        "Kintone由来のため、この欄を手で書き換えても翌日の自動連携で元に戻ります。"
        "直したいときは Kintone 側を修正してください。", warn=True)

# ════════════════════════════════════════════════
# 7. 項目 4・5
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "項目くわしく ④⑤　住まい・福祉用具・請求", eyebrow="FIELDS", bar=C_WELFARE)
add_text(s, "④ 入居施設名・居室番号・在宅区分", 0.7, 1.8, 11, 0.4, size=15, bold=True, color=C_WELFARE)
add_table(s, [
    ["項目", "入れる内容", "選べる値"],
    ["入居施設名 🕐", "入居先の施設名（在宅は空欄でOK）", "文字入力"],
    ["居室番号 🕐", "施設の部屋番号", "例：101"],
    ["在宅区分 🕐", "住まいの種類", "ー／自宅／外部施設／その他"],
], 0.7, 2.2, 11.9, 1.5, [2.7, 5.8, 3.4], header_color=C_WELFARE)
add_text(s, "⑤ 福祉用具利用者・レセプトチェック対象・請求区分", 0.7, 3.95, 11, 0.4, size=15, bold=True, color=C_WELFARE)
add_table(s, [
    ["項目", "入れる内容", "選べる値"],
    ["福祉用具利用者 🕐", "福祉用具を使う方はチェック（介護保険・自費・販売すべて含む）", "チェック ✓"],
    ["レセプトチェック対象 🕐", "ふだんは未設定（自動判定）でOK", "✓オン=強制追加／□オフ=強制除外／未設定=自動"],
    ["請求区分 🕐", "請求の種類を選択", "ー／自費レンタル／介護保険レンタル／併用"],
], 0.7, 4.35, 11.9, 2.0, [2.7, 5.8, 3.4], header_color=C_WELFARE)

# ════════════════════════════════════════════════
# 8. 項目 6・7
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "項目くわしく ⑥⑦　ケアマネ・介護保険", eyebrow="FIELDS", bar=C_CM)
add_text(s, "⑥ ケアマネージャー情報", 0.7, 1.8, 11, 0.4, size=15, bold=True, color=C_CM)
add_table(s, [
    ["項目", "入れる内容"],
    ["居宅介護支援事業所 🕐", "担当ケアマネの所属事業所名"],
    ["担当CM 🕐", "担当ケアマネジャーのお名前"],
], 0.7, 2.2, 11.9, 1.2, [3.6, 8.3], header_color=C_CM)
add_text(s, "⑦ 介護保険情報", 0.7, 3.6, 11, 0.4, size=15, bold=True, color=C_ACCENT)
add_table(s, [
    ["項目", "入れる内容", "選べる値"],
    ["要介護度 🕐", "認定区分", "ー／申請中／要支援1・2／要介護1〜5"],
    ["負担割合 🕐", "利用者の自己負担割合", "ー／1割／2割／3割"],
    ["介護保険被保険者証 🕐", "確認できているか", "ー／確認済／未確認"],
    ["介護保険負担割合証 🕐", "確認できているか", "ー／確認済／未確認"],
], 0.7, 4.0, 11.9, 2.0, [3.0, 4.6, 4.3], header_color=C_ACCENT)
add_text(s, "※「ー」は「まだ未設定」の意味。わかった時点で選び直してください。",
         0.7, 6.15, 11.5, 0.4, size=12, color=C_MUTED)

# ════════════════════════════════════════════════
# 9. 項目 8・9・10
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "項目くわしく ⑧⑨⑩　支払い・連絡先・カイポケ", eyebrow="FIELDS")
add_table(s, [
    ["項目", "入れる内容", "選べる値"],
    ["⑧ 支払い区分 🕐", "生活保護受給かどうか（レセプトの生保判定に使用）", "ー／非生保／生保"],
    ["⑨ キーパーソン：氏名", "ご家族など連絡先になる方の氏名", "文字入力"],
    ["⑨ キーパーソン：続柄", "本人との関係（長男・妻 など）", "文字入力"],
    ["⑨ キーパーソン：連絡先", "電話番号など", "文字入力"],
    ["⑩ カイポケ登録（基本情報）", "カイポケに基本情報を登録済みか", "未登録／登録済"],
], 0.7, 2.1, 11.9, 3.0, [3.4, 5.5, 3.0])

# ════════════════════════════════════════════════
# 10. 変更履歴 🕐
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "🕐 変更履歴の使い方", eyebrow="HISTORY")
add_lines(s, [
    ("「いつから介護度が変わったか」など、変化の履歴を残せます。", False, C_INK),
    ("過去の月を正しく表示するために役立ちます。", False, C_MUTED),
], 0.7, 1.8, 11.9, 0.9, size=15)
add_text(s, "● 自動で記録される", 0.7, 2.7, 11, 0.4, size=15, bold=True, color=C_ACCENT)
add_text(s, "編集モードで🕐つきの項目を変えると「いつから有効ですか？」と日付をきく画面が出ます。日付を入れると履歴に残ります。",
         0.9, 3.1, 11.5, 0.7, size=13, color=C_INK)
add_text(s, "● 過去の履歴を手で追加する", 0.7, 3.85, 11, 0.4, size=15, bold=True, color=C_ACCENT)
steps2 = [
    ("1", "編集モードにする", "右上の「編集」を押す"),
    ("2", "🕐 を押す", "項目名の横の🕐で履歴の画面を開く"),
    ("3", "「履歴を追加」", "値・実効日（必須）・備考を入れて追加 →「保存する」"),
]
x = 0.7
for n, h, b in steps2:
    add_rect(s, x, 4.3, 3.75, 1.5, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, n, x + 0.2, 4.4, 1.0, 0.7, size=30, bold=True, color=C_ACCENT)
    add_text(s, h, x + 0.2, 5.0, 3.4, 0.4, size=14, bold=True, color=C_INK)
    add_text(s, b, x + 0.2, 5.35, 3.4, 0.5, size=11, color=C_MUTED)
    x += 4.05
callout(s, 0.7, 6.0, 11.9, 0.85, "🗑️ まちがえたら削除できます",
        "同じ🕐の画面で「削除」を押せば消せます。この履歴は施設契約情報とちがい、消えずに保存されます。")

# ════════════════════════════════════════════════
# 11. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="Q&A")
qa = [
    ("入力したのに反映されません", "自動保存はありません。「編集」中に右上の「保存する」を押したか確認してください。うっかり保存せず移動しようとしても確認ダイアログが出るので大丈夫です。"),
    ("施設契約情報を直したのに翌日もとに戻る", "Kintoneからの自動連携情報です。直すときは Kintone 側を修正してください。"),
    ("「ー」のままでいい？", "「ー」は未設定の意味。わかる範囲で選ぶと、レセプトチェックや集計が正しく動きます。"),
    ("レンタルの新規・解約はどこ？", "「利用者新規・変更情報入力」タブの「レンタル契約情報」です。基本情報タブの施設契約情報は施設の入退去専用。"),
]
y = 1.95
for q, a in qa:
    add_text(s, "Q. " + q, 0.7, y, 11.9, 0.45, size=15, bold=True, color=C_ACCENT)
    add_text(s, "A. " + a, 1.0, y + 0.42, 11.4, 0.7, size=13, color=C_INK)
    y += 1.2

out = "docs/基本情報タブ_操作マニュアル.pptx"
prs.save(out)
print("作成しました:", out, "／ スライド数:", len(prs.slides._sldIdLst))
