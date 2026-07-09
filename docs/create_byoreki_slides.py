"""
WelfareAssist Pro - 病歴・状態タブ 操作マニュアル（スライド生成）
出力: docs/病歴・状態タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_byoreki_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（他タブマニュアルと共通）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 病歴・状態タブ固有色
C_MED      = RGBColor(0xDC, 0x26, 0x26)   # 赤: 病歴セクション
C_MEDLT    = RGBColor(0xFE, 0xE2, 0xE2)
C_AI       = RGBColor(0x7C, 0x3A, 0xED)   # 紫: AI用具提案
C_AILT     = RGBColor(0xED, 0xE9, 0xFE)
C_OCR      = RGBColor(0x1D, 0x4E, 0xD8)   # 青: OCR文書読み取り
C_OCRLT    = RGBColor(0xDB, 0xEA, 0xFE)
C_PLAN     = RGBColor(0xD9, 0x77, 0x06)   # 黄: 選定予定用具
C_PLANLT   = RGBColor(0xFF, 0xF7, 0xED)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_MED):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro  |  病歴・状態タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def callout(slide, l, t, w, h, title, body, warn=False, ai=False, ocr=False, plan=False):
    if ai:
        bg, edge = C_AILT, C_AI
    elif ocr:
        bg, edge = C_OCRLT, C_OCR
    elif plan:
        bg, edge = C_PLANLT, C_PLAN
    elif warn:
        bg, edge = C_WARNLT, C_WARN
    else:
        bg, edge = C_ACCENTLT, C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_MED, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_MEDLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_MED)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5,
         size=16, bold=True, color=C_MED)
add_text(s, "「病歴・状態」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=46, bold=True, color=C_INK)
add_rect(s, 1.05, 3.95, 2.2, 0.08, fill=C_MED)
add_lines(s, [
    ("利用者の病名・身体状況を記録する画面です。", False, C_MUTED),
    ("診療情報提供書などの医療文書をAIが読み取り、病歴欄へ自動入力することもできます。", False, C_MUTED),
    ("また病歴をもとにAIが適切な福祉用具を提案します。", False, C_MUTED),
], 1.05, 4.25, 11, 1.2, size=17)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：利用者をクリック → 「病歴・状態」タブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_MED)

# ════════════════════════════════════════════════
# 2. このタブでできること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")

feats = [
    (C_MEDLT, C_MED, "📋", "病歴・身体状況を記録",
     "病名・麻痺の有無・ADLなど利用者の状態を自由テキストで入力・保存します。全タブから参照されます。"),
    (C_OCRLT, C_OCR, "📄", "医療文書をAI読み取り",
     "診療情報提供書・退院サマリーなどのPDF・画像をアップロードするとAIがテキストを抽出します。"),
    (C_AILT, C_AI, "🤖", "病歴から用具を提案",
     "入力済みの病歴をもとにAIが適切な福祉用具の種目を提案します。選定時の参考にできます。"),
    (C_PLANLT, C_PLAN, "📦", "選定予定の用具を登録",
     "品名・種目を仮登録できます。「福祉用具選定」タブの正式選定前のメモとして活用できます。"),
]
x = 0.7
for bg, edge, ic, title, body in feats:
    add_rect(s, x, 2.0, 2.8, 3.4, fill=bg, line=edge, lw=1)
    add_rect(s, x, 2.0, 0.1, 3.4, fill=edge)
    add_text(s, ic, x + 0.25, 2.15, 0.8, 0.7, size=26, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x + 0.22, 2.9, 2.45, 0.55, size=15, bold=True, color=C_INK)
    add_text(s, body, x + 0.22, 3.5, 2.45, 1.65, size=12, color=C_MUTED)
    x += 2.98

callout(s, 0.7, 5.7, 11.9, 0.95,
        "自動保存について",
        "入力内容は1.2秒後に自動的に保存されます。「保存」ボタンを押さなくても大丈夫です。ただし医療文書の読み取り結果は「病歴欄に反映」ボタンを押すまで保存されません。")

# ════════════════════════════════════════════════
# 3. 病歴・身体状況を入力する
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "病歴・身体状況を入力する", eyebrow="MEDICAL HISTORY")

# テキストエリアのモックアップ
add_rect(s, 0.7, 2.0, 11.9, 0.42, fill=C_MEDLT)
add_rect(s, 0.7, 2.0, 0.08, 0.42, fill=C_MED)
add_text(s, "病歴・身体状況", 0.95, 2.06, 5.0, 0.3, size=14, bold=True, color=C_MED)
add_rect(s, 0.7, 2.42, 11.9, 2.45, fill=C_WHITE, line=C_LINE, lw=1)
add_lines(s, [
    ("【病名】脳梗塞（右片麻痺）、高血圧症、糖尿病", False, C_INK),
    ("【麻痺】右上下肢に軽度麻痺あり。握力低下。", False, C_INK),
    ("【ADL】歩行は伝い歩き可。階段は要介助。", False, C_INK),
    ("【認知】日常会話は問題なし。新しいことの記憶に難あり。", False, C_MUTED),
    ("【その他】退院後1ヶ月。リハビリ継続中。", False, C_MUTED),
], 0.95, 2.52, 11.2, 2.25, size=13)

add_table(s, [
    ["入力項目（参考）", "記入内容の例"],
    ["病名", "脳梗塞、骨折、パーキンソン病、心疾患 など"],
    ["麻痺の有無", "右片麻痺、下肢麻痺なし など"],
    ["ADL（日常生活動作）", "歩行・立ち座り・入浴・排泄の状態"],
    ["認知機能", "日常会話可能、要確認 など"],
    ["その他特記事項", "退院後の経過・リハビリ状況・家族の状況 など"],
], 0.7, 5.1, 11.9, 2.12, [2.6, 9.3], size=13)

# ════════════════════════════════════════════════
# 4. 医療文書をAIで読み取る
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "医療文書をAIで読み取る", eyebrow="OCR / AI DOCUMENT READING", bar=C_OCR)

steps = [
    ("1", "文書エリアに\nドロップ or クリック",
     "テキストエリアの下にあるアップロードエリアをクリック（またはファイルをドラッグ）します。"),
    ("2", "医療文書を選択",
     "診療情報提供書・退院サマリーなどのPDF・PNG・JPG・WEBPファイルを選択します。"),
    ("3", "AIがテキストを抽出",
     "「AI解析中...」と表示されます。完了すると読み取り結果が画面に表示されます。"),
    ("4", "「病歴欄に反映」を押す",
     "結果を確認したら「病歴欄に反映」ボタンを押します。病歴テキストエリアに内容が追記されます。"),
]
x = 0.7
for n, h, b in steps:
    add_rect(s, x, 2.1, 2.82, 3.0, fill=C_WHITE, line=C_OCRLT, lw=1.5)
    add_rect(s, x, 2.1, 0.09, 3.0, fill=C_OCR)
    add_rect(s, x + 0.25, 2.22, 0.55, 0.55, fill=C_OCR)
    add_text(s, n, x + 0.25, 2.22, 0.55, 0.55, size=20, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, h, x + 0.25, 2.88, 2.38, 0.7, size=14, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.65, 2.38, 1.25, size=12, color=C_MUTED)
    x += 3.05

callout(s, 0.7, 5.32, 11.9, 0.82,
        "対応ファイル形式",
        "PDF・PNG・JPG・JPEG・WEBP（最大20MB）。診療情報提供書・退院サマリー・入院証明書などに対応しています。",
        ocr=True)

callout(s, 0.7, 6.28, 11.9, 0.82,
        "⚠ 「病歴欄に反映」を押すまで保存されません",
        "読み取り結果は画面表示のみです。病歴テキストに反映するには「病歴欄に反映」ボタンを押してください。反映後は1.2秒後に自動保存されます。",
        warn=True)

# ════════════════════════════════════════════════
# 5. 病歴から用具を提案（AI）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "病歴から福祉用具を提案（AI）", eyebrow="AI EQUIPMENT SUGGESTION", bar=C_AI)

add_rect(s, 0.7, 2.0, 11.9, 1.1, fill=C_AILT)
add_rect(s, 0.7, 2.0, 0.1, 1.1, fill=C_AI)
add_lines(s, [
    ("病歴・身体状況テキストをもとに、AIが適切な福祉用具の種目を提案する機能です。", False, C_INK),
    ("画面右上の「病歴から用具を提案」ボタン（紫）を押すと、AI分析が始まります（数秒〜10秒）。", False, C_INK),
], 0.95, 2.1, 11.2, 0.88, size=14)

# 提案結果のモックアップ
add_rect(s, 0.7, 3.3, 11.9, 0.38, fill=C_AILT)
add_rect(s, 0.7, 3.3, 0.08, 0.38, fill=C_AI)
add_text(s, "AIによる提案", 0.95, 3.36, 5.0, 0.26, size=13, bold=True, color=C_AI)
add_rect(s, 0.7, 3.68, 11.9, 2.15, fill=RGBColor(0xF5, 0xF3, 0xFF), line=C_AILT, lw=1)
add_lines(s, [
    ("【特殊寝台・特殊寝台付属品】", True, C_AI),
    ("右片麻痺・歩行不安定のため、離床・就寝動作の補助に特殊寝台（電動2/3モーター）と", False, C_INK),
    ("サイドレール・介助バーをお勧めします。", False, C_INK),
    ("", False, C_INK),
    ("【歩行補助器（歩行器）】", True, C_AI),
    ("室内の伝い歩き・廊下移動の安全確保のため、歩行器（交互型）が適合すると考えられます。", False, C_INK),
], 0.95, 3.75, 11.1, 1.95, size=13)

callout(s, 0.7, 6.02, 11.9, 0.72,
        "提案内容は参考情報です",
        "AI提案は診断や処方ではありません。実際の選定は利用者・家族・ケアマネとの相談のうえで行ってください。",
        warn=True)

callout(s, 0.7, 6.85, 11.9, 0.42,
        "病歴テキストが空の場合は提案できません",
        "先に病歴・身体状況テキストを入力してからご利用ください。",
        ai=True)

# ════════════════════════════════════════════════
# 6. 選定予定の福祉用具
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "選定予定の福祉用具を登録する", eyebrow="PLANNED EQUIPMENT", bar=C_PLAN)

add_rect(s, 0.7, 2.0, 11.9, 1.05, fill=C_PLANLT)
add_rect(s, 0.7, 2.0, 0.1, 1.05, fill=C_PLAN)
add_lines(s, [
    ("病歴・状態タブの下部に「選定予定の福祉用具」セクションがあります。", False, C_INK),
    ("「福祉用具選定」タブで正式に機器を決める前の仮メモとして品名・種目を登録しておけます。", False, C_INK),
], 0.95, 2.1, 11.2, 0.88, size=14)

# モックアップ: 選定予定セクション
add_rect(s, 0.7, 3.22, 11.9, 0.42, fill=RGBColor(0xFE, 0xF9, 0xC3))
add_text(s, "選定予定の福祉用具", 0.95, 3.3, 5.0, 0.26, size=13, bold=True, color=C_PLAN)
add_text(s, "+ 追加", 11.5, 3.3, 1.0, 0.26, size=11, bold=True, color=C_MUTED, align=PP_ALIGN.RIGHT)

rows_mock = [
    ("特殊寝台", "特殊寝台"),
    ("介助バー", "特殊寝台付属品"),
    ("歩行器（屋内用）", "歩行器"),
]
y = 3.72
for fname, cat in rows_mock:
    add_rect(s, 0.7, y, 11.9, 0.48, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.9, y + 0.1, 4.5, 0.28, fill=C_LINE)
    add_text(s, fname, 0.95, y + 0.11, 4.4, 0.26, size=12, color=C_INK)
    add_rect(s, 5.55, y + 0.1, 2.8, 0.28, fill=C_LINE)
    add_text(s, cat, 5.6, y + 0.11, 2.75, 0.26, size=12, color=C_MUTED)
    y += 0.56

add_table(s, [
    ["操作", "方法"],
    ["項目を追加する", "「＋ 追加」ボタンをクリック → 品名・種目を入力"],
    ["項目を編集する", "直接テキストをクリックして書き換える（編集モード時）"],
    ["項目を削除する", "行右端の「×」ボタンをクリック"],
    ["正式な機器登録に変換", "「福祉用具選定」タブで改めて詳細情報を入力してください"],
], 0.7, 5.6, 11.9, 1.65, [3.0, 8.9], size=13)

# ════════════════════════════════════════════════
# 7. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="FAQ")

qas = [
    ("病歴はどのくらい詳しく書けばいいですか？",
     "AI用具提案や議事録生成の精度を上げるため、病名・麻痺の有無・ADL状態・認知機能・退院状況などできるだけ詳しく記載することをお勧めします。"),
    ("医療文書のアップロードができません",
     "対応形式（PDF・PNG・JPG・JPEG・WEBP）かご確認ください。また20MBを超えるファイルはアップロードできません。"),
    ("AI読み取り結果が一部おかしいです",
     "AIの読み取り精度はドキュメントの品質によって異なります。「病歴欄に反映」後に内容を確認・修正してください。"),
    ("「病歴から用具を提案」ボタンが押せません（グレーになっている）",
     "病歴テキストエリアに内容が入力されていない場合は提案できません。先に病歴を入力してください。"),
    ("選定予定の用具は「福祉用具選定」タブに自動連携されますか？",
     "されません。選定予定はあくまでメモです。「福祉用具選定」タブで改めて正式な機器情報を入力してください。"),
]
y = 2.0
rh = 0.96
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_MED)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_MED)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.38, size=14, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.44, 10.9, 0.44, size=12, color=C_MUTED)
    y += rh + 0.08

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "病歴・状態タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
