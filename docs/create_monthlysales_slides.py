"""
WelfareAssist Pro - 月次売上処理タブ 操作マニュアル（スライド生成）
出力: docs/月次売上処理タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_monthlysales_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（他タブマニュアルと共通のベース＋種類別色）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# 種類別色（アプリ画面のタブ色と統一）
C_INS      = RGBColor(0x25, 0x63, 0xEB)   # 青: 介護保険レンタル
C_INSLT    = RGBColor(0xDB, 0xEA, 0xFE)
C_SELF     = RGBColor(0x7C, 0x3A, 0xED)   # 紫: 自費レンタル
C_SELFLT   = RGBColor(0xED, 0xE9, 0xFE)
C_SALE     = RGBColor(0x16, 0xA3, 0x4A)   # 緑: 販売
C_SALELT   = RGBColor(0xDC, 0xFC, 0xE7)
C_AMBER    = RGBColor(0xB4, 0x53, 0x09)
C_AMBERLT  = RGBColor(0xFE, 0xF3, 0xC7)
C_DELETE   = RGBColor(0xDC, 0x26, 0x26)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=28, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.6, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro ｜ 月次売上処理タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 1 else C_ACCENTLT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def callout(slide, l, t, w, h, title, body, warn=False, accentcolor=None):
    if warn:
        bg, edge = C_WARNLT, C_WARN
    elif accentcolor:
        bg, edge = accentcolor
    else:
        bg, edge = C_ACCENTLT, C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


def step_cards(slide, steps, y=2.1, h=2.6, gap=0.32):
    n = len(steps)
    total_w = 11.9
    w = (total_w - gap * (n - 1)) / n
    x = 0.7
    for n_label, head, body in steps:
        add_rect(slide, x, y, w, h, fill=C_WHITE, line=C_LINE, lw=1)
        add_text(slide, n_label, x + 0.22, y + 0.12, 1.0, 0.85, size=36, bold=True, color=C_ACCENT)
        add_text(slide, head, x + 0.22, y + 0.95, w - 0.44, 0.55, size=15, bold=True, color=C_INK)
        add_text(slide, body, x + 0.22, y + 1.5, w - 0.44, h - 1.6, size=11.5, color=C_MUTED)
        x += w + gap


# ════════════════════════════════════════════════
# 1. タイトル
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 2.0, 11, 0.5, size=16, bold=True, color=C_ACCENT)
add_text(s, "「月次売上処理」タブの使い方", 1.0, 2.6, 11.3, 1.3, size=44, bold=True, color=C_INK)
add_rect(s, 1.05, 3.9, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("毎月の売上（介護保険レンタル・自費レンタル・販売）を確認し、", False, C_MUTED),
    ("カイポケCSVを取り込み、金額を確定してCSVを出力する画面です。", False, C_MUTED),
    ("ここで確定した金額が「売上・仕入突合」タブの基準値になります。", False, C_MUTED),
], 1.05, 4.2, 11, 1.4, size=16)
add_text(s, "対象：福祉用具専門相談員のみなさま　／　場所：左サイドバー「月次売上処理」（紫ボタン）",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════════════════════════════════════════════
# 2. できること
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")
feats = [
    (C_ACCENTLT, C_ACCENT, "📋", "売上データの確認",
     "選んだ月・事業所の介護保険レンタル/自費レンタル/販売の一覧と合計金額を確認できます。"),
    (C_INSLT, C_INS, "📥", "カイポケCSVの取込み",
     "介護保険レンタルタブでカイポケの2つのCSVを取り込み、当月分のデータに丸ごと入れ替えます。"),
    (C_SALELT, RGBColor(0x16,0x65,0x34), "🔒", "金額の確定",
     "3種類の売上をそれぞれ確定すると、その時点の金額がスナップショットとして固定されます。"),
]
x = 0.7
for bg, edge, ic, title, body in feats:
    add_rect(s, x, 2.0, 3.8, 3.1, fill=bg, line=edge, lw=1)
    add_rect(s, x, 2.0, 0.1, 3.1, fill=edge)
    add_text(s, ic, x + 0.3, 2.2, 0.8, 0.7, size=28, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x + 0.28, 2.95, 3.3, 0.55, size=17, bold=True, color=C_INK)
    add_text(s, body, x + 0.28, 3.55, 3.35, 1.4, size=12.5, color=C_MUTED)
    x += 4.07
callout(s, 0.7, 5.4, 11.9, 1.0,
        "📝 この画面で確定した金額が「売上・仕入突合」タブの基準になります",
        "月末の作業は、まずこの画面で3種類（介護保険レンタル・自費レンタル・販売）すべてを確定させることから始めます。")

# ════════════════════════════════════════════════
# 3. 画面の基本構成
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "画面の基本構成", eyebrow="LAYOUT")
add_text(s, "月度・事業所を選び、3つのタブで種類ごとの一覧を切り替えます。", 0.7, 1.95, 11.9, 0.4, size=14, color=C_MUTED)

# タブ風ボタン
add_rect(s, 0.7, 2.5, 3.6, 0.55, fill=C_INS)
add_text(s, "🏥 介護保険レンタル 312件", 0.7, 2.5, 3.6, 0.55, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 4.4, 2.5, 3.0, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "💰 自費レンタル 48件", 4.4, 2.5, 3.0, 0.55, size=13, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 7.5, 2.5, 2.6, 0.55, fill=RGBColor(0xE5,0xE7,0xEB))
add_text(s, "🛒 販売 23件", 7.5, 2.5, 2.6, 0.55, size=13, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_table(s, [
    ["タブ", "データの出どころ", "集計の基準日"],
    ["介護保険レンタル", "このタブ内のカイポケCSVインポートで取り込んだデータ", "利用開始日〜終了日が対象月と重なる品目"],
    ["自費レンタル", "基本情報タブ（ClientDetail）の入力データ", "利用開始日〜終了日が対象月と重なる品目"],
    ["販売", "基本情報タブ（ClientDetail）の入力データ", "納品日が対象月内の品目のみ"],
], 0.7, 3.3, 11.9, 1.9, [2.5, 5.4, 4.0], size=13)

callout(s, 0.7, 5.5, 11.9, 1.0,
        "⚠ 事業所が未設定の利用者は「全事業所」でも表示されません",
        "事業所フィルターは利用者ごとの「事業所」設定で判定されます。売上が出ない利用者がいたら、まず基本情報タブで事業所設定を確認してください。",
        warn=True)

# ════════════════════════════════════════════════
# 4. 売上サマリーの見方
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "売上サマリーの見方", eyebrow="SUMMARY")
add_text(s, "タブの下に常に表示されている4枚のカードです（税抜の金額）。", 0.7, 1.9, 11.9, 0.4, size=14, color=C_MUTED)

cards = [
    (C_INSLT, C_INS, "介護保険レンタル", "312件", "¥4,560,000", "確定", False),
    (C_SELFLT, C_SELF, "自費レンタル", "48件", "¥612,000", "解除", True),
    (C_SALELT, RGBColor(0x16,0x65,0x34), "販売", "23件", "¥481,300", "確定", False),
    (RGBColor(0xF1,0xF5,0xF1), C_INK, "合計", "383件", "¥5,653,300", None, None),
]
x = 0.7
for bg, fg, label, cnt, amt, btn, confirmed in cards:
    add_rect(s, x, 2.4, 2.85, 2.0, fill=bg)
    add_text(s, label, x + 0.2, 2.55, 2.5, 0.35, size=12.5, bold=True, color=fg)
    if confirmed:
        add_text(s, "✓ 確定済", x + 0.2, 2.88, 2.5, 0.3, size=10, bold=True, color=RGBColor(0x16,0x65,0x34))
    add_text(s, cnt, x + 0.2, 3.15, 2.5, 0.45, size=20, bold=True, color=fg)
    add_text(s, amt, x + 0.2, 3.6, 2.5, 0.35, size=13, color=fg)
    if btn:
        btncol = C_DELETE if btn == "解除" else fg
        add_rect(s, x + 0.2, 4.05, 2.45, 0.32, fill=C_WHITE if btn == "解除" else fg, line=btncol, lw=0.75)
        add_text(s, btn, x + 0.2, 4.05, 2.45, 0.32, size=11, bold=True,
                 color=btncol if btn == "解除" else C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += 3.07

add_table(s, [
    ["状態", "内容"],
    ["未確定（「確定」ボタン）", "件数・金額はその場で自動計算された最新の値。元データを直せば数字も変わる。"],
    ["確定済（「解除」ボタン）", "確定した瞬間の件数・金額をスナップショットとして固定。以後、元データを変えても変わらない。"],
    ["解除できないとき", "「売上・仕入突合」タブで月次確定が完了している場合、先にそちらを解除する必要がある。"],
], 0.7, 4.75, 11.9, 1.85, [3.5, 8.4], size=12.5)

# ════════════════════════════════════════════════
# 5. カイポケCSVインポート：ファイル選択
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "カイポケCSVインポート（1）ファイル選択", eyebrow="IMPORT STEP 1", bar=C_INS)
add_text(s, "介護保険レンタルのデータは月に1回、カイポケCSVを取り込みます。取り込むと当月分が丸ごと入れ替わります（洗い替え）。",
         0.7, 1.9, 11.9, 0.6, size=13.5, color=C_MUTED)

y = 2.65
for num, name, hint, must_col in [
    ("1", "① サービスチェックシート.csv　*必須", "用具の種類・単位数などの明細データ", C_DELETE),
    ("2", "② 利用者請求.csv　*必須", "ないと請求額が0件になります", C_DELETE),
]:
    add_rect(s, 0.7, y, 11.9, 0.85, fill=RGBColor(0xF9,0xFA,0xFB), line=C_LINE, lw=0.5)
    add_rect(s, 0.85, y + 0.18, 0.5, 0.5, fill=must_col)
    add_text(s, num, 0.85, y + 0.18, 0.5, 0.5, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, name, 1.55, y + 0.12, 8.0, 0.35, size=14, bold=True, color=C_INK)
    add_text(s, hint, 1.55, y + 0.46, 8.0, 0.3, size=11.5, color=C_MUTED)
    add_rect(s, 9.9, y + 0.2, 1.8, 0.45, fill=C_INSLT)
    add_text(s, "ファイルを選択", 9.9, y + 0.2, 1.8, 0.45, size=11, bold=True, color=C_INS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.0

# ボタン列
bx = 0.7
for label, bg, fg in [("🔍 プレビュー", C_INSLT, C_INS), ("📥 インポート実行", C_INS, C_WHITE), ("✕ 選択クリア", RGBColor(0xF1,0xF5,0xF9), RGBColor(0x47,0x55,0x69))]:
    add_rect(s, bx, y + 0.15, 2.1, 0.5, fill=bg)
    add_text(s, label, bx, y + 0.15, 2.1, 0.5, size=12, bold=True, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += 2.25
add_rect(s, 9.9, y + 0.15, 2.0, 0.5, fill=C_DELETE)
add_text(s, "🗑 データクリア", 9.9, y + 0.15, 2.0, 0.5, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

callout(s, 0.7, y + 0.9, 11.9, 0.95,
        "⚠ 2つのファイルは必ず両方選んでください",
        "「②利用者請求.csv」を選ばずに取り込むと、用具は登録されますが請求額（給付対象金額）がすべて0円になります。",
        warn=True)

# ════════════════════════════════════════════════
# 6. プレビュー確認
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "カイポケCSVインポート（2）プレビュー確認", eyebrow="IMPORT STEP 2", bar=C_INS)
add_text(s, "①のファイルを選ぶと「プレビュー」が押せます。実際に取り込む前に必ず確認しましょう。", 0.7, 1.9, 11.9, 0.4, size=13.5, color=C_MUTED)

boxes = [("マッチ成功", "302名", RGBColor(0x16,0x65,0x34)), ("未マッチ", "3名", C_DELETE),
         ("品目数", "340件", C_INS), ("給付対象金額", "¥4,560,000", C_INS)]
x = 0.7
for label, val, col in boxes:
    add_rect(s, x, 2.5, 2.85, 1.3, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, label, x + 0.2, 2.65, 2.5, 0.3, size=11.5, color=C_MUTED)
    add_text(s, val, x + 0.2, 3.0, 2.5, 0.6, size=20, bold=True, color=col)
    x += 3.07

callout(s, 0.7, 4.1, 11.9, 1.1,
        "👀 「未マッチ」が0件でなければ「詳細表示」を確認",
        "カイポケ側の氏名・カナが利用者マスターと違う場合に発生しやすいです。未マッチの利用者はそのままインポートには含まれません。")

callout(s, 0.7, 5.35, 11.9, 1.35,
        "🖱️ 請求データの紐づけは「クリック→クリック」だけ",
        "「請求金額の紐づけ状況」で未紐づけがあれば「詳細表示」。左側で未紐づけ利用者をクリックして選択（青くなる）→右側の対応する請求データをクリックすると緑色の「紐づけ済み」になります。「インポート実行」を押すまで何度でもやり直せます。")

# ════════════════════════════════════════════════
# 7. インポート実行と警告
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "カイポケCSVインポート（3）実行と警告確認", eyebrow="IMPORT STEP 3", bar=C_INS)
add_text(s, "プレビューの内容に問題がなければ「インポート実行」を押します。完了すると成功メッセージが表示されます。",
         0.7, 1.9, 11.9, 0.5, size=13.5, color=C_MUTED)

add_rect(s, 0.7, 2.6, 11.9, 0.6, fill=C_SALELT)
add_text(s, "✓ インポート完了：302名の利用者に340件の用具を登録しました", 0.9, 2.6, 11.5, 0.6,
         size=13, bold=True, color=RGBColor(0x16,0x65,0x34), anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, 0.7, 3.4, 11.9, 1.0, fill=C_AMBERLT)
add_rect(s, 0.7, 3.4, 0.08, 1.0, fill=C_AMBER)
add_lines(s, [
    ("⚠️ 請求データ未紐づけ：3名（売上集計から除外されています）", True, C_AMBER),
    ("対象者の給付対象金額は売上合計に含まれていません。確定前に必ず確認・再インポートしてください。", False, C_INK),
], 0.95, 3.55, 11.4, 0.85, size=13)

callout(s, 0.7, 4.6, 11.9, 1.5,
        "🗑 データクリアボタンについて",
        "全事業所・全利用者の介護保険レンタルデータを一括削除するボタンです。確認ダイアログが出ますが、この操作は取り消せません。自費レンタル・販売のデータには影響しません。通常の月次作業では使わず、データ異常時のリセット用です。",
        warn=True)

# ════════════════════════════════════════════════
# 8. 販売タブの金額計算
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "販売タブの金額計算（重要）", eyebrow="CALCULATION", bar=C_SALE)
add_text(s, "販売の合計金額は、いくつかの項目を順番に足し引きして「総計」を出します。", 0.7, 1.9, 11.9, 0.4, size=13.5, color=C_MUTED)

flow = ["小計\n単価×数量", "消費税\n小計×税率", "税込金額\n小計+消費税", "+送料+送料消費税\n+調整額", "総計\n負担額の基準"]
x = 0.7
w = 2.15
for i, item in enumerate(flow):
    head, sub = item.split("\n")
    is_last = (i == len(flow) - 1)
    add_rect(s, x, 2.5, w, 1.1, fill=C_SALELT if is_last else C_WHITE, line=C_LINE, lw=1)
    add_text(s, head, x + 0.1, 2.62, w - 0.2, 0.4, size=13, bold=True,
             color=RGBColor(0x16,0x65,0x34) if is_last else C_INK, align=PP_ALIGN.CENTER)
    add_text(s, sub, x + 0.1, 3.05, w - 0.2, 0.45, size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
    x += w
    if not is_last:
        add_text(s, "→", x, 2.85, 0.3, 0.4, size=16, color=C_MUTED, align=PP_ALIGN.CENTER)
        x += 0.3

add_table(s, [
    ["列", "計算式", "備考"],
    ["小計", "単価 × 数量", ""],
    ["消費税", "小計 × 税率（切り捨て）", "税区分：10%/軽8%/非課税/税込。税込商品は0円"],
    ["税込金額", "小計 + 消費税", ""],
    ["送料（税抜）／送料消費税", "手動入力 ／ 送料×10%（四捨五入）", "送料0円なら消費税も0円"],
    ["調整額", "手動入力（±）", "端数調整・割引などに使用"],
    ["総計", "税込金額 + 送料（税抜）+ 送料消費税 + 調整額", "利用者負担額・申請額の計算基準"],
], 0.7, 4.1, 11.9, 2.6, [3.2, 5.0, 3.7], size=11.5)

# ════════════════════════════════════════════════
# 9. 利用者負担額の計算ルール
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "利用者負担額・申請額の自動計算", eyebrow="BURDEN CALC", bar=C_SALE)
add_text(s, "「利用者自己負担割合」を設定すると、利用者負担額と申請額が自動で決まります（手入力済みの場合はそちらが優先）。",
         0.7, 1.9, 11.9, 0.6, size=13.5, color=C_MUTED)

add_table(s, [
    ["利用者自己負担割合", "利用者負担額", "申請額"],
    ["自己負担０（日常生活給付）", "0円", "総計"],
    ["一部負担（日常生活給付）", "上限額（総計との小さい方）", "総計 − 利用者負担額"],
    ["1〜3割負担（受領委任払い）", "総計×割合（切り上げ／上限額あればその範囲内）", "総計 − 利用者負担額"],
    ["全額負担（償還払い）", "総計", "総計"],
], 0.7, 2.7, 11.9, 2.0, [4.0, 4.4, 3.5], size=13)

callout(s, 0.7, 5.0, 11.9, 1.1,
        "📌 「一部負担」「1〜3割負担」では上限額を必ず入力",
        "「一部負担時の上限額」が0円のままだと、利用者負担額が正しく計算されません。")

# ════════════════════════════════════════════════
# 10. CSV出力
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "CSV出力", eyebrow="EXPORT")
add_text(s, "各タブの右上「CSVダウンロード」ボタンで、表示中の月・事業所のデータをそのまま出力できます。",
         0.7, 1.9, 11.9, 0.4, size=13.5, color=C_MUTED)

add_table(s, [
    ["タブ", "主な列"],
    ["介護保険レンタル", "あおぞらID・氏名・施設名・商品名・種類・メーカー・卸会社・単位数・タイスコード・利用開始/終了日"],
    ["自費レンタル", "あおぞらID・氏名・施設名・商品名・単価・個数・金額（税抜/税込）・税区分・利用開始/終了日・取引方法・備考"],
    ["販売", "金額内訳（小計〜総計）＋受注日・納品日・支払方法・自己負担割合・上限額・負担額・申請額・申請市町村・営業担当・備考ほか全26列"],
], 0.7, 2.6, 11.9, 2.7, [2.5, 9.4], size=12.5)

# ════════════════════════════════════════════════
# 11. 毎月のチェックリスト
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "毎月の作業チェックリスト", eyebrow="CHECKLIST")
items = [
    "介護保険レンタルタブでカイポケCSV（2ファイル）をインポート済み",
    "インポート後に「請求データ未紐づけ」の警告バナーが出ていないか確認した",
    "自費レンタル・販売の一覧が最新か確認した（前日の自動更新後）",
    "事業所フィルターを正しく設定してCSVを出力した",
    "3種類（介護保険レンタル・自費レンタル・販売）をすべて確定した",
    "（このあと）「売上・仕入突合」タブで請求書と突合し、月次確定まで進める",
]
y = 2.0
for it in items:
    add_rect(s, 0.7, y, 11.9, 0.62, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.9, y + 0.16, 0.3, 0.3, fill=None, line=C_ACCENT, lw=1.5)
    add_text(s, it, 1.4, y, 11.0, 0.62, size=13.5, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.72

# ════════════════════════════════════════════════
# 12. よくある質問
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="FAQ")
qas = [
    ("CSVを取り込んだら前月分のデータが消えました",
     "介護保険レンタルのインポートは選択中の月度に対する洗い替え方式です。前月分は前月の月度を選んで確認してください。"),
    ("確定後に元データを修正したら金額は変わりますか？",
     "変わりません。確定はスナップショットです。反映させるには一度「解除」して再度「確定」してください。"),
    ("「解除」ボタンが押せません",
     "「売上・仕入突合」タブで月次確定まで完了している場合、そちらを先に解除する必要があります。"),
    ("自費レンタル・販売の件数がここに出てきません",
     "対象利用者の事業所設定、利用開始/終了日（自費レンタル）、納品日（販売）が対象月の範囲に入っているか確認してください。"),
    ("給付対象金額が0円の利用者がいます",
     "「利用者請求.csv」に紐づいていない可能性があります。インポートのプレビューで紐づけ状況を確認し、必要なら再インポートしてください。"),
]
y = 2.0
rh = 0.98
for q, a in qas:
    add_rect(s, 0.7, y, 11.9, rh, fill=C_WHITE, line=C_LINE, lw=0.5)
    add_rect(s, 0.7, y, 0.06, rh, fill=C_ACCENT)
    add_text(s, "Q", 0.9, y + 0.1, 0.5, rh - 0.2, size=15, bold=True, color=C_ACCENT)
    add_text(s, q, 1.3, y + 0.05, 10.9, 0.4, size=13.5, bold=True, color=C_INK)
    add_text(s, a, 1.3, y + 0.46, 10.9, 0.48, size=11.5, color=C_MUTED)
    y += rh + 0.08

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "月次売上処理タブ_操作マニュアル.pptx")
prs.save(out)
print("saved: " + out)
