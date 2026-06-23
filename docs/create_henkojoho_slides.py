"""
WelfareAssist Pro - 利用者新規・変更情報入力タブ 操作マニュアル（スライド生成）
出力: docs/利用者新規・変更情報入力タブ_操作マニュアル.pptx

Googleドライブにアップロードすると Google スライドに自動変換されます。
実行: python docs/create_henkojoho_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（基本情報マニュアルと共通）──
C_GROUND   = RGBColor(0xEA, 0xF1, 0xED)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_INK      = RGBColor(0x1C, 0x2B, 0x27)
C_MUTED    = RGBColor(0x5B, 0x6E, 0x68)
C_ACCENT   = RGBColor(0x0E, 0x8A, 0x78)
C_ACCENTLT = RGBColor(0xE2, 0xF0, 0xEC)
C_WARN     = RGBColor(0xE2, 0x71, 0x4A)
C_WARNLT   = RGBColor(0xFB, 0xEA, 0xE2)
C_LINE     = RGBColor(0xD7, 0xE2, 0xDC)
# カード種別の色（画面と対応）
C_PENDING  = RGBColor(0xD9, 0x77, 0x06)
C_HOSP     = RGBColor(0xEA, 0x58, 0x0C)
C_CHANGE   = RGBColor(0x05, 0x96, 0x69)
C_CONTRACT = RGBColor(0x7C, 0x3A, 0xED)
C_GRAY     = RGBColor(0x94, 0xA3, 0xB8)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return box


def add_lines(slide, lines, l, t, w, h, size=16, spacing=1.25, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return box


def content_layout(slide, title, eyebrow=None, bar=C_ACCENT):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
    add_rect(slide, 0, 0, 0.28, 7.5, fill=bar)
    if eyebrow:
        add_text(slide, eyebrow, 0.7, 0.42, 11, 0.4, size=13, bold=True, color=bar)
    add_text(slide, title, 0.66, 0.72, 12, 0.9, size=30, bold=True, color=C_INK)
    add_rect(slide, 0.7, 1.62, 1.4, 0.06, fill=bar)
    add_text(slide, "WelfareAssist Pro ｜ 利用者新規・変更情報入力タブ 操作マニュアル",
             0.7, 7.05, 12, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, rows, l, t, w, h, col_w, header_color=C_ACCENT, size=13):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if ri == 0 else (C_WHITE if ri % 2 == 1 else C_ACCENTLT)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_INK
    return gt


def callout(slide, l, t, w, h, title, body, warn=False):
    bg = C_WARNLT if warn else C_ACCENTLT
    edge = C_WARN if warn else C_ACCENT
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.08, h, fill=edge)
    add_lines(slide, [(title, True, edge), (body, False, C_INK)],
              l + 0.25, t + 0.12, w - 0.4, h - 0.2, size=13)


# ════════ 1. タイトル ════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_GROUND)
add_rect(s, 0, 0, 13.33, 0.22, fill=C_ACCENT)
add_text(s, "福祉用具マネージャー　操作マニュアル", 1.0, 1.9, 11, 0.5, size=16, bold=True, color=C_ACCENT)
add_text(s, "「利用者新規・変更情報入力」タブの使い方", 1.0, 2.5, 11.5, 1.5, size=38, bold=True, color=C_INK)
add_rect(s, 1.05, 4.0, 2.2, 0.08, fill=C_ACCENT)
add_lines(s, [
    ("利用者の「いつ・何が変わったか」を記録する画面です。", False, C_MUTED),
    ("新規契約・入院・退院・解約・変更などを、やさしくまとめました。", False, C_MUTED),
], 1.05, 4.3, 11, 1.2, size=17)
add_text(s, "場所：利用者をクリック →「利用者新規・変更情報入力」タブ",
         1.05, 6.4, 11.5, 0.5, size=13, bold=True, color=C_ACCENT)

# ════════ 2. できること ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "このタブでできること", eyebrow="OVERVIEW")
add_rect(s, 0.7, 2.0, 11.9, 1.6, fill=C_ACCENTLT)
add_rect(s, 0.7, 2.0, 0.08, 1.6, fill=C_ACCENT)
add_lines(s, [
    ("このタブは利用者の「できごとの記録ノート」です。", True, C_INK),
    ("1つのできごと＝1枚のカードとして、種類（情報種別）と日付を記録します。", False, C_INK),
    ("入院でサービスを止めた／退院で再開した／解約した／新しく契約した——", False, C_INK),
    ("こうした変化をここに残すと、レセプトチェックやスプレッドシートに反映されます。", False, C_INK),
], 1.0, 2.2, 11.4, 1.4, size=15)
add_text(s, "記録できるできごと（情報種別）", 0.7, 4.0, 11, 0.45, size=16, bold=True, color=C_ACCENT)
add_table(s, [
    ["分類", "種別"],
    ["契約", "新規 ／ 解約"],
    ["入退院", "入院（サービス停止）／ 退院（サービス再開）"],
    ["その他", "変更あり ／ その他 ／ デモ"],
], 0.7, 4.45, 11.9, 1.9, [2.4, 9.5])

# ════════ 3. 3ステップ ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "入力の3ステップ", eyebrow="BASIC")
steps = [
    ("1", "「＋情報を追加」", "右上のボタンで、黄色い「入力中」カードが最上部に出る。"),
    ("2", "種別と内容を入力", "情報種別を選び、日付・利用区分・記録者・特記を入力。"),
    ("3", "「保存する」", "右上の保存で確定。スプレッドシートにも自動反映。"),
]
x = 0.7
for n, h, b in steps:
    add_rect(s, x, 2.1, 3.75, 2.4, fill=C_WHITE, line=C_LINE, lw=1)
    add_text(s, n, x + 0.25, 2.25, 1.2, 1.0, size=44, bold=True, color=C_ACCENT)
    add_text(s, h, x + 0.25, 3.25, 3.3, 0.5, size=18, bold=True, color=C_INK)
    add_text(s, b, x + 0.25, 3.75, 3.3, 0.9, size=13, color=C_MUTED)
    x += 4.05
callout(s, 0.7, 4.85, 11.9, 1.0,
        "📤 保存すると変更情報スプレッドシートへ自動同期（数秒後）",
        "手動でシートを触る必要はありません。記録は毎日の自動更新後も残ります"
        "（Kintone自動連携の記録を除く）。")

# ════════ 4. 情報種別 ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "情報種別（カードの種類）", eyebrow="TYPES")
add_text(s, "「何のできごとか」を選びます。種別ごとに入力する日付が変わります。",
         0.7, 1.8, 11.9, 0.4, size=14, color=C_MUTED)
add_table(s, [
    ["情報種別", "使う場面", "入力する日付"],
    ["新規", "福祉用具レンタルを新しく開始した", "請求開始日（新規）"],
    ["入院", "入院でサービスを停止した", "請求停止日（入院）"],
    ["退院", "退院でサービスを再開した", "請求開始日（退院）"],
    ["解約", "レンタルを解約した", "請求停止日（解約）"],
    ["変更あり", "用具の変更など（請求に直接関わらない）", "—（特記に記入）"],
    ["その他", "上記にあてはまらない連絡・メモ", "—（特記に記入）"],
    ["デモ", "デモ機の貸し出し", "デモ開始日・デモ終了日"],
], 0.7, 2.3, 11.9, 3.5, [2.0, 6.4, 3.5])
add_text(s, "※ 施設への入居・退去（施設契約情報）は「基本情報」タブに移動しました。",
         0.7, 6.0, 11.5, 0.4, size=12, color=C_WARN)

# ════════ 5. 利用区分 ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "利用区分（チェックボックス）", eyebrow="CATEGORY")
add_lines(s, [
    ("各カードに「利用区分」のチェックボックスがあります（複数選べます）。", False, C_INK),
    ("チェックした内容はスプレッドシートの「利用区分」列（G列）に出力されます。", False, C_INK),
], 0.7, 1.9, 11.9, 0.9, size=15)
opts = ["☑ 自費レンタル", "☑ 介護保険レンタル", "☑ 販売"]
x = 0.7
for o in opts:
    add_rect(s, x, 3.0, 3.75, 0.9, fill=C_ACCENTLT, line=C_ACCENT, lw=1)
    add_text(s, o, x, 3.0, 3.75, 0.9, size=16, bold=True, color=C_INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += 4.05
callout(s, 0.7, 4.3, 11.9, 1.0,
        "💡 新しいカードは初め「介護保険レンタル」にチェック",
        "違う場合はチェックを外し、あてはまるものを選んでください（複数可）。")

# ════════ 6. 画面の並び ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "画面の並び（カードの表示順）", eyebrow="MAP")
rows_map = [
    ("1", "入力中", "「＋情報を追加」直後。保存まで最上部に固定", C_PENDING, "黄色"),
    ("2", "入院・退院情報", "入院→退院を日付でペア表示（Kintone自動連携）", C_HOSP, "オレンジ"),
    ("3", "変更あり・その他・デモ", "請求に直接関わらない変化・メモ・デモ", C_CHANGE, "緑/水色"),
    ("4", "レンタル契約情報", "新規→解約をペア表示（福祉用具レンタル契約）", C_CONTRACT, "紫"),
    ("5", "解約（単独）", "新規とペアにならなかった解約", C_GRAY, "グレー"),
]
y = 2.0
rh = 0.78
for n, title, desc, col, badge in rows_map:
    add_rect(s, 0.7, y, 0.5, rh - 0.12, fill=col)
    add_text(s, n, 0.7, y - 0.06, 0.5, rh, size=18, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, 1.35, y - 0.06, 6.0, rh, size=16, bold=True, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, desc, 7.4, y - 0.06, 4.4, rh, size=12, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, badge, 11.9, y - 0.06, 1.0, rh, size=12, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    y += rh + 0.05

# ════════ 7. ペアリング ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "ペアリングの仕組み", eyebrow="PAIRING")
add_text(s, "関係する2枚のカードは、自動で左右に並べて表示されます。",
         0.7, 1.8, 11.9, 0.4, size=14, color=C_MUTED)
add_rect(s, 0.7, 2.4, 11.9, 1.7, fill=C_WHITE, line=C_LINE, lw=1)
add_text(s, "入院 ↔ 退院", 1.0, 2.55, 11, 0.45, size=16, bold=True, color=C_HOSP)
add_lines(s, [
    ("入院カードと、その入院日以降のいちばん早い退院カードが自動でペアになります。", False, C_INK),
    ("退院がまだなら「退院情報なし」と表示されます。", False, C_MUTED),
], 1.0, 3.05, 11.3, 1.0, size=13)
add_rect(s, 0.7, 4.3, 11.9, 1.9, fill=C_WHITE, line=C_LINE, lw=1)
add_text(s, "新規 ↔ 解約", 1.0, 4.45, 11, 0.45, size=16, bold=True, color=C_CONTRACT)
add_lines(s, [
    ("日付で自動ペアされます。正しくない場合は、解約カード右上のドロップダウンから", False, C_INK),
    ("ペア先の新規を手動で選び直せます（手動指定が優先・橙色の「手動ペア」表示）。", False, C_INK),
], 1.0, 4.95, 11.3, 1.1, size=13)

# ════════ 8. 自動 vs 手動 ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "自動で入る記録・手で入れる記録", eyebrow="SOURCE")
add_table(s, [
    ["種類", "入りかた", "編集"],
    ["入院・退院", "Kintone（入退院情報）から自動連携", "原則そのまま。手で直しても翌日の同期で戻る"],
    ["新規・解約・変更あり・その他・デモ", "このタブで手入力", "自由に追加・編集・削除でき、保持される"],
], 0.7, 2.1, 11.9, 1.7, [3.6, 4.2, 4.1])
callout(s, 0.7, 4.2, 11.9, 1.7,
        "⚠️ Kintone自動連携の記録（入院・退院・施設入居）はアプリで書き換えても翌日戻る",
        "直したいときは Kintone 側を修正してください。"
        "手入力で追加した記録（新規・解約など）はそのまま保持されます。", warn=True)

# ════════ 9. スプレッドシート反映 ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "スプレッドシートへの反映", eyebrow="SHEET")
add_text(s, "保存すると、変更情報スプレッドシート（共有）へ自動で追記されます。利用区分はG列に出力されます。",
         0.7, 1.9, 11.9, 0.7, size=15, color=C_INK)
callout(s, 0.7, 2.8, 11.9, 1.7,
        "📌 シートは「追記専用」です",
        "新しく作った記録は新しい行として追記されますが、すでにシートにある過去の行は、"
        "あとから利用区分などを変えても更新されません。利用区分は新しく入力する記録から反映されます。", warn=True)
callout(s, 0.7, 4.7, 11.9, 1.3,
        "🔀 シートの行はヘッダーごと並べ替えないでください",
        "ヘッダーを含めた全選択ソートは表示崩れの原因になります。並べ替えはフィルター表示を使ってください。")

# ════════ 10. Q&A ════════
s = prs.slides.add_slide(BLANK)
content_layout(s, "よくある質問", eyebrow="Q&A")
qa = [
    ("追加したのに保存できない／消えた", "黄色い「入力中」のままかも。右上の「保存する」で確定。キャンセルで破棄。"),
    ("入院・退院を直したのに翌日戻る", "Kintone自動連携のため上書きされます。Kintone側で直してください。"),
    ("解約が違う新規とペアになる", "解約カード右上のドロップダウンから正しい新規を選び直せます（手動優先）。"),
    ("利用区分がシートに出ない", "保存後に自動同期。ただし過去に出力済みの行は更新されません（追記専用）。"),
    ("施設の入居・退去はどこ？", "「基本情報」タブの「施設契約情報」です。"),
]
y = 1.95
for q, a in qa:
    add_text(s, "Q. " + q, 0.7, y, 11.9, 0.4, size=14, bold=True, color=C_ACCENT)
    add_text(s, "A. " + a, 1.0, y + 0.36, 11.4, 0.6, size=12, color=C_INK)
    y += 1.02

out = "docs/利用者新規・変更情報入力タブ_操作マニュアル.pptx"
prs.save(out)
print("作成しました:", out, "／ スライド数:", len(prs.slides._sldIdLst))
